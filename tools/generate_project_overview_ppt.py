#!/usr/bin/env python3
"""生成 QEMU-VCS CoSim Platform 项目总览 PPT"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# 颜色
DARK_BLUE = RGBColor(0x00, 0x2B, 0x5C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
HEADER_BG = RGBColor(0x00, 0x3D, 0x7A)
ROW_ALT = RGBColor(0xE8, 0xF0, 0xFA)

CN_FONT = "SimHei"


def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BLUE
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(8.4), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = CN_FONT
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = subtitle
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(0xA0, 0xC0, 0xE0)
    p2.font.name = CN_FONT
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(20)


def add_content_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0), Inches(0), Inches(10), Inches(0.9))
    tf = title_box.text_frame
    tf.paragraphs[0].text = title
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = DARK_BLUE
    tf.paragraphs[0].font.name = CN_FONT
    tf.paragraphs[0].alignment = PP_ALIGN.LEFT
    tf.margin_left = Inches(0.5)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    body_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(9), Inches(6))
    tf2 = body_box.text_frame
    tf2.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.text = b
        p.font.size = Pt(15)
        p.font.color.rgb = DARK_GRAY
        p.font.name = CN_FONT
        p.space_before = Pt(6)


def add_table_slide(prs, title, headers, rows):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0), Inches(0), Inches(10), Inches(0.9))
    tf = title_box.text_frame
    tf.paragraphs[0].text = title
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = DARK_BLUE
    tf.paragraphs[0].font.name = CN_FONT
    tf.margin_left = Inches(0.5)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    n_rows = len(rows) + 1
    n_cols = len(headers)
    width = Inches(9.2)
    table_shape = slide.shapes.add_table(n_rows, n_cols, Inches(0.4), Inches(1.2), width, Inches(0.4) * n_rows)
    table = table_shape.table
    col_w = int(width / n_cols)
    for ci in range(n_cols):
        table.columns[ci].width = col_w
    for ci, h in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = CN_FONT
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = HEADER_BG
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12)
                p.font.color.rgb = BLACK
                p.font.name = CN_FONT
                p.alignment = PP_ALIGN.LEFT
            if ri % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = ROW_ALT


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ── 1. 封面 ──
    add_title_slide(prs,
        "QEMU-VCS CoSim Platform 项目总览",
        "协同仿真平台 · 架构设计 · 部署指南\n2026-04-23")

    # ── 2. 项目简介 ──
    add_content_slide(prs, "项目简介",
        [
            "• QEMU-VCS RTL 级软硬件协同仿真",
            "• PCIe TLP 级 MMIO/DMA/MSI",
            "• 支持 SHM 本地和 TCP 跨机两种模式",
            "• 支持 UVM VIP 验证 IP",
        ])

    # ── 3. 系统架构 ──
    add_content_slide(prs, "系统架构",
        [
            "QEMU 侧:",
            "  Guest Linux + virtio-net driver → cosim-pcie-rc 自定义设备 → libcosim_bridge.so",
            "",
            "VCS 侧:",
            "  simv_vip (UVM + pcie_tl_vip) → bridge_vcs DPI-C → ETH SHM → eth_tap_bridge → TAP",
            "",
            "通信:",
            "  SHM (POSIX 共享内存 + Unix Socket) 或 TCP (v2 三连接 ctrl+data+aux)",
        ])

    # ── 4. 核心模块 ──
    add_table_slide(prs, "核心模块",
        ["模块", "路径", "功能"],
        [
            ["bridge/common", "bridge/common/", "共享库(SHM/ring buffer/transport/trace)"],
            ["bridge/qemu", "bridge/qemu/", "QEMU 侧桥接 + irq_poller"],
            ["bridge/vcs", "bridge/vcs/", "VCS 侧桥接 + DPI-C + virtqueue DMA"],
            ["bridge/eth", "bridge/eth/", "以太网层 (ETH SHM + MAC DPI)"],
            ["qemu-plugin", "qemu-plugin/", "cosim-pcie-rc PCIe RC 设备"],
            ["vcs-tb", "vcs-tb/", "VCS testbench (EP stub + glue + VIP top)"],
            ["pcie_tl_vip", "pcie_tl_vip/", "UVM PCIe TL 验证 IP"],
        ])

    # ── 5. 编译工具链 ──
    add_table_slide(prs, "编译工具链",
        ["工具", "版本要求", "用途"],
        [
            ["gcc/g++", ">= 4.8", "Bridge C 编译"],
            ["cmake", ">= 3.16", "Bridge CMake 构建"],
            ["python3", ">= 3.8", "QEMU meson 构建"],
            ["meson + ninja", "latest", "QEMU 编译"],
            ["VCS", "Q-2020+", "RTL 仿真"],
            ["glib", ">= 2.66", "QEMU 依赖(可自动源码编译)"],
        ])

    # ── 6. 安装流程 ──
    add_content_slide(prs, "安装流程",
        [
            "三种部署模式:",
            "  • ./setup.sh --mode local: 同机全栈(QEMU+VCS+TAP)",
            "  • ./setup.sh --mode qemu-only: QEMU 侧远程",
            "  • ./setup.sh --mode vcs-only: VCS 侧远程",
            "",
            "setup.sh 自动执行流程:",
            "  1. 依赖检测",
            "  2. Bridge 编译",
            "  3. QEMU 编译",
            "  4. VCS VIP 编译",
            "  5. TAP 编译 + setcap",
            "  6. 单元测试",
        ])

    # ── 7. 部署模式对比 ──
    add_table_slide(prs, "部署模式对比",
        ["特性", "Local SHM", "TCP 跨机"],
        [
            ["通信方式", "POSIX SHM+Socket", "TCP v2 三连接"],
            ["QEMU 命令", "--shm /cosim0 --sock /tmp/cosim0.sock", "--transport tcp --port-base 9100"],
            ["VCS 命令", "--shm --sock", "--transport tcp --remote-host IP"],
            ["串口交互", "--serial-sock PATH", "--serial-sock PATH"],
            ["波形 dump", "默认 FSDB", "默认 FSDB"],
            ["仿真速度", "较快(SHM 零拷贝)", "较慢(TCP 延迟)"],
        ])

    # ── 8. TCP Transport 架构 ──
    add_content_slide(prs, "TCP Transport 架构",
        [
            "v2 三连接:",
            "  • ctrl (port_base+0) + data (+1) + aux (+2)",
            "",
            "各 channel 职责:",
            "  • ctrl: sync 消息 (TLP_READY, CPL_READY, DMA_CPL, SHUTDOWN)",
            "  • data: TLP/CPL 数据",
            "  • aux: DMA_REQ, DMA_CPL, DMA_DATA, MSI, ETH",
            "",
            "关键机制:",
            "  • irq_poller: QEMU 侧独立线程, MSG_PEEK 按类型分发",
            "  • 自适应 poll 超时: 1ms → 5ms → 50ms",
        ])

    # ── 9. VIP 验证 IP ──
    add_content_slide(prs, "VIP 验证 IP",
        [
            "pcie_tl_vip: 完整 UVM PCIe TL VIP",
            "",
            "  • Agent: RC/EP driver + monitor + scoreboard",
            "  • Sequence: MRd/MWr/CfgRd/CfgWr/DMA/MSI/Error injection",
            "  • cosim_rc_driver: DPI-C polling + handle_completion 回传",
            "  • cosim_vip_top: NOTIFY event + RX poll + MSI 注入 + FSDB dump",
        ])

    # ── 10. 功能测试 ──
    add_content_slide(prs, "功能测试",
        [
            "• ./cosim.sh test phase1-5: 自动编排测试",
            "• ./cosim.sh test-guide: 交互式测试向导(ping/iperf/arping/压力)",
            "• --serial-sock: Guest 串口交互(python/socat)",
            "• 波形: 默认 FSDB dump, +NO_WAVE 关闭",
        ])

    # ── 11. 验证结果 ──
    add_table_slide(prs, "验证结果",
        ["测试项", "Local SHM", "TCP 跨机"],
        [
            ["PCIe TLP", "✅ 168+", "✅ 100K+"],
            ["DMA read/write", "✅ 280", "✅ 4319r+2197w"],
            ["MSI 中断", "✅", "✅ 10万次"],
            ["VQ-TX→TAP", "✅ 47包", "✅ 1066包"],
            ["波形 FSDB", "✅", "✅"],
            ["tag error", "0", "0"],
        ])

    # ── 12. 用户扩展指南 ──
    add_content_slide(prs, "用户扩展指南",
        [
            "• 自定义 VCS testbench: 修改 vcs-tb/*.sv, make vcs-vip 重编",
            "• 自定义 EP 行为: 修改 pcie_ep_stub.sv 的寄存器和 completion 逻辑",
            "• 新增 DPI-C 函数: 在 bridge/vcs/bridge_vcs.c 添加, bridge_vcs.sv 声明 import",
            "• 自定义 Guest: buildroot menuconfig 或替换 rootfs",
            "• 新增测试 sequence: 在 pcie_tl_vip/src/seq/ 添加, cosim_test.sv 引用",
            "• 编译命令: make vcs-vip (VIP模式), make bridge (仅Bridge), cmake + ninja (QEMU)",
        ])

    # ── 13. 注意事项 ──
    add_content_slide(prs, "注意事项",
        [
            "• QEMU 源码树需同步最新 bridge 代码(setup.sh 自动处理)",
            "• eth_tap_bridge 每次编译后需 sudo setcap cap_net_admin+ep",
            "• VCS 机器需 source ~/set-env.sh 加载 EDA 环境",
            "• Guest 建议用 buildroot rootfs(含 virtio_net 驱动)",
            "• 仿真速度由 VCS RTL 仿真决定, ping 超时设 600s+",
            "• TCP 启动顺序: 先 QEMU(listen) → 再 VCS(connect) → 再 TAP",
        ])

    out = "/home/ubuntu/ryan/software/cosim-platform/docs/cosim_project_overview.pptx"
    prs.save(out)
    print(f"PPT 已生成: {out}")


if __name__ == "__main__":
    main()
