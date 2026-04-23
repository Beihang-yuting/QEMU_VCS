#!/usr/bin/env python3
"""生成 CoSim Platform 问题排查与修复报告 PPT"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# 颜色
DARK_BLUE = RGBColor(0x00, 0x2B, 0x5C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
HEADER_BG = RGBColor(0x00, 0x3D, 0x7A)
ROW_ALT = RGBColor(0xE8, 0xF0, 0xFA)

CN_FONT = "SimHei"


def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BLUE
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(8.4), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = CN_FONT
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = subtitle
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(0xA0, 0xC0, 0xE0)
    p2.font.name = CN_FONT
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(20)


def add_content_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0), Inches(0), Inches(10), Inches(0.9))
    tf = title_box.text_frame
    tf.paragraphs[0].text = title
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = DARK_BLUE
    tf.paragraphs[0].font.name = CN_FONT
    tf.paragraphs[0].alignment = PP_ALIGN.LEFT
    tf.margin_left = Inches(0.5)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    body_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(9), Inches(6))
    tf2 = body_box.text_frame
    tf2.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.text = b
        p.font.size = Pt(15)
        p.font.color.rgb = DARK_GRAY
        p.font.name = CN_FONT
        p.space_before = Pt(6)


def add_table_slide(prs, title, headers, rows):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0), Inches(0), Inches(10), Inches(0.9))
    tf = title_box.text_frame
    tf.paragraphs[0].text = title
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = DARK_BLUE
    tf.paragraphs[0].font.name = CN_FONT
    tf.margin_left = Inches(0.5)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    n_rows = len(rows) + 1
    n_cols = len(headers)
    width = Inches(9.2)
    table_shape = slide.shapes.add_table(n_rows, n_cols, Inches(0.4), Inches(1.2), width, Inches(0.4) * n_rows)
    table = table_shape.table
    col_w = int(width / n_cols)
    for ci in range(n_cols):
        table.columns[ci].width = col_w
    for ci, h in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = CN_FONT
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = HEADER_BG
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12)
                p.font.color.rgb = BLACK
                p.font.name = CN_FONT
                p.alignment = PP_ALIGN.LEFT
            if ri % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = ROW_ALT


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ── 封面 ──
    add_title_slide(prs,
        "CoSim Platform 问题排查与修复报告",
        "Setup 流程集成 · VIP MSI 修复 · TCP 消息交错分析\n2026-04-23")

    # ── 问题总览 ──
    add_table_slide(prs, "问题总览",
        ["#", "问题", "严重性", "状态"],
        [
            ["1", "cosim.sh 不支持 TCP transport 参数", "高", "✅ 已修复"],
            ["2", "cosim.sh 不支持磁盘镜像启动 (--drive)", "高", "✅ 已修复"],
            ["3", "setup.sh NEED_TAP_BRIDGE 归属错误", "高", "✅ 已修复"],
            ["4", "setup.sh VCS 编译为 legacy 而非 VIP", "严重", "✅ 已修复"],
            ["5", "setup.sh VCS 源文件列表缺失 transport_tcp.c 等", "严重", "✅ 已修复"],
            ["6", "cosim_vip_top.sv VIP 模式缺少 MSI 中断注入", "严重", "✅ 已修复"],
            ["7", "glue_if_to_stub.sv stub_isr_set 多驱动冲突", "高", "✅ 已修复"],
            ["8", "TCP transport DMA read 消费 TLP_READY 导致协议错位", "严重", "✅ 已修复"],
            ["9", "NOTIFY 在 always_ff 中调阻塞 DPI-C 冻结仿真", "严重", "✅ 已修复"],
            ["10", "_nb 函数消费不匹配消息导致 aux channel 协议错位", "严重", "✅ 已修复"],
            ["11", "dma_write_bytes 先发 DMA_DATA 堵住 aux channel", "严重", "✅ 已修复"],
        ])

    # ── 问题 1-2 ──
    add_content_slide(prs, "问题 1-2: cosim.sh 不支持 TCP / 磁盘镜像",
        [
            "现象: setup.sh 提示的跨机命令无法执行",
            "  ./cosim.sh start qemu --transport tcp --port-base 9100 → 报错「未知选项」",
            "",
            "根因: cmd_start_qemu 只有 --shm/--sock/--initrd 三个参数",
            "  -device 参数硬编码为 SHM 模式",
            "",
            "修复: 新增 --transport, --port-base, --instance-id, --drive, --append",
            "  TCP: -device cosim-pcie-rc,transport=tcp,port_base=N,instance_id=N",
            "  Drive: -drive file=X,format=raw/qcow2,if=virtio",
            "",
            "VCS 侧同步修复: cmd_start_vcs 新增 --transport tcp --remote-host IP",
            "  TCP 必须指定 --remote-host, 缺失时报错提示",
        ])

    # ── 问题 3 ──
    add_content_slide(prs, "问题 3: NEED_TAP_BRIDGE 归属错误",
        [
            "现象: vcs-only 模式不编译 eth_tap_bridge",
            "",
            "跨机拓扑: eth_tap_bridge 运行在 VCS 机器 (61) 上",
            "  VCS 机器: simv_vip → ETH SHM → eth_tap_bridge → TAP cosim0",
            "  QEMU 机器: 只有 QEMU + Guest, 不需要 TAP",
            "",
            "原代码:   qemu-only: NEED_TAP_BRIDGE=true ← 错误",
            "          vcs-only:  NEED_TAP_BRIDGE=false ← 错误",
            "",
            "修复: 翻转归属",
            "  qemu-only: NEED_TAP_BRIDGE=false",
            "  vcs-only:  NEED_TAP_BRIDGE=true",
        ])

    # ── 问题 4-5 ──
    add_content_slide(prs, "问题 4-5: setup.sh VCS 编译模式错误",
        [
            "现象: setup.sh 编译产出 legacy simv, 跨机联调需要 VIP simv_vip",
            "",
            "缺失项:",
            "  +define+COSIM_VIP_MODE — VIP 宏开关",
            "  -ntb_opts uvm-1.2 — UVM 框架",
            "  pcie_tl_if.sv / pcie_tl_pkg.sv — VIP 源文件",
            "  transport_tcp.c / transport_shm.c / trace_log.c — C 源文件",
            "  sock_sync.c 错误引用（应为 sock_sync_vcs.c）",
            "",
            "修复: setup.sh 直接调用 make vcs-vip",
            "  产出路径统一: build/simv_vip",
            "  cosim.sh resolve_simv 优先搜 build/simv_vip",
        ])

    # ── 问题 6-7 ──
    add_content_slide(prs, "问题 6-7: VIP 模式 MSI 中断未注入",
        [
            "现象: Guest virtio0 中断 = 0, TX packets = 0",
            "  virtio_net driver 未完成初始化, TX queue 未激活",
            "",
            "根因: cosim_vip_top.sv RX poll 成功后缺少 MSI 触发",
            "  tb_top.sv (legacy): isr_set → bridge_vcs_raise_msi(0)  ✅",
            "  cosim_vip_top.sv (VIP): 无 MSI 代码  ❌",
            "  virtqueue_dma.c 注释: \"tb_top.sv handles it\" — 但 VIP 不走 tb_top",
            "",
            "修复:",
            "  cosim_vip_top.sv: RX inject 后 stub_isr_set 脉冲 + raise_msi",
            "  glue_if_to_stub.sv: stub_isr_set 从 output 改为 input",
            "  时序关键: 先设 ISR bit → 再发 MSI（Guest 读 ISR 必须非零）",
        ])

    # ── 问题 8 ──
    add_content_slide(prs, "问题 8: TCP DMA Read TLP 缓存",
        [
            "现象: DMA read 期间 QEMU 发的 TLP_READY 被 recv_sync 消费",
            "  pending 计数与 data channel 消息错位",
            "",
            "修复: TLP ring buffer 缓存（1024 entry）",
            "  DMA read 遇到 TLP_READY 时立即 recv_tlp 并 push 到缓存",
            "  poll_tlp 优先从缓存 pop",
            "  自适应超时 1ms→5ms→50ms 适配跨机 TCP 延迟",
        ])

    # ── 问题 9 ──
    add_content_slide(prs, "问题 9: NOTIFY always_ff 冻结仿真",
        [
            "现象: handle_vio_notify 在 always_ff 中调阻塞 DPI-C",
            "  bridge_dma_read_bytes recv_sync 阻塞 → 整个仿真冻结",
            "  TLP 缓存溢出 (cache full)",
            "",
            "修复: NOTIFY 处理从 always_ff 移到 initial 块",
            "  always_ff 只做边沿检测并触发 SV event（非阻塞）",
            "  initial 块等待 event 后调 handle_vio_notify（可安全阻塞）",
        ])

    # ── 问题 10 ──
    add_content_slide(prs, "问题 10: _nb 函数 aux channel 协议错位",
        [
            "现象: recv_dma_req_nb 读 header 发现非 DMA_REQ 返回 -1",
            "  但 header 已消费, payload 留在 buffer → 协议错位",
            "  irq_poller 永远读不到有效消息",
            "",
            "修复: _nb 函数改用 MSG_PEEK 预读 header",
            "  类型不匹配时返回 1（无数据）不消费 buffer",
            "  匹配时正式 recv 消费 header + payload",
        ])

    # ── 问题 11（关键突破）──
    add_content_slide(prs, "问题 11: DMA write 顺序导致 aux channel 僵死 [关键突破]",
        [
            "现象: dma_write_bytes 先发 DMA_DATA 再发 DMA_REQ",
            "  irq_poller MSG_PEEK 看到 DMA_DATA（非 DMA_REQ）→ 不匹配 → 返回 1",
            "  DMA_REQ 在 DMA_DATA 后面永远读不到",
            "  → aux channel 僵死 → ctrl recv_sync 断开 → 所有后续 DMA 失败",
            "",
            "修复: 调换顺序 — 先 send_dma_req 再 send_dma_data",
            "  irq_poller PEEK 到 DMA_REQ → 匹配 → cosim_dma_cb → recv_dma_data",
            "",
            "效果: DMA read/write 全部正常, MSI 中断注入成功, 数据面打通",
        ])

    # ── TCP channel 架构 ──
    add_content_slide(prs, "TCP Transport Channel 架构分析",
        [
            "当前 v2 三连接架构:",
            "  ctrl (port_base+0): sync_msg — TLP_READY, CPL_READY, DMA_CPL, SHUTDOWN",
            "  data (port_base+1): TLP, CPL 数据",
            "  aux  (port_base+2): DMA_REQ, DMA_CPL, MSI, ETH, DMA_DATA",
            "",
            "问题: ctrl channel 是共享的",
            "  TLP_READY 和 DMA_CPL 都走 ctrl channel",
            "  DMA read 循环 recv_sync 消费本属于 poll 的 TLP_READY",
            "  消费后只计数 (pending_tlp++), 不读 data channel 上的 TLP 数据",
            "  后续 recv_tlp 与 data channel 上的消息错位",
            "",
            "解决方案 (选项):",
            "  A: DMA read 消费 TLP_READY 时同步 recv_tlp 缓存完整 TLP",
            "  B: 将 TLP_READY 和 DMA_CPL 拆分到不同 channel",
            "  C: 消息自描述 — data channel 每条消息带 type header",
        ])

    # ── 验证结果 ──
    add_table_slide(prs, "跨机测试验证结果 (53 ↔ 61)",
        ["测试项", "结果", "说明"],
        [
            ["QEMU TCP listen (53)", "✅", "端口 9100 v2 三连接握手"],
            ["VCS TCP connect (61)", "✅", "ctrl+data+aux 正确配对"],
            ["ETH SHM + TAP bridge", "✅", "/cosim_eth0 + cosim0 10.0.0.1"],
            ["Guest boot + eth0", "✅", "buildroot + virtio_net 10.0.0.2"],
            ["PCIe TLP 交换", "✅", "33000+ TLP 持续流动"],
            ["Completion 回传", "✅", "30700+ CPL via handle_completion"],
            ["DMA read (TCP)", "✅", "48 次成功"],
            ["DMA write (TCP)", "✅", "51 次成功"],
            ["MSI 中断注入", "✅", "12 次 pci_set_irq"],
            ["VQ-TX → ETH SHM → TAP", "✅", "4 包转发成功"],
            ["TAP → ETH SHM → VQ-RX", "✅", "11 包注入成功"],
            ["Virtio 数据面", "✅", "双向数据面完全打通"],
        ])

    # ── 提交记录 ──
    add_content_slide(prs, "修复提交记录 (共 8 个 commit)",
        [
            "1. feat(setup): cosim.sh TCP/drive 参数 + setup.sh TAP/VCS 修正",
            "2. fix(vcs): make vcs-vip + MSI 注入 + resolve_simv 路径",
            "3. fix(glue): stub_isr_set output→input 消除多驱动冲突",
            "4. fix(vip): handle_completion 覆盖回传 QEMU",
            "5. fix(tcp): TLP ring buffer 缓存 + 自适应 poll 超时",
            "6. fix(vip): NOTIFY always_ff→initial event 解除仿真冻结",
            "7. fix(tcp): _nb 函数 MSG_PEEK 避免 aux 协议错位",
            "8. fix(tcp): dma_write_bytes 先 DMA_REQ 再 DMA_DATA [关键突破]",
        ])

    # ── 结论 ──
    add_content_slide(prs, "结论: Virtio 数据面完全打通",
        [
            "跨机 TCP 模式 (53 QEMU ↔ 61 VCS) 端到端数据面已打通:",
            "",
            "  Guest TX: virtio_net → NOTIFY → VQ-TX → DMA read → ETH SHM → TAP",
            "  Guest RX: TAP → ETH SHM → VQ-RX → DMA write → MSI → virtio_net",
            "",
            "验证数据 (持续运行中):",
            "  TLP: 33000+  |  Completion: 30700+  |  NOTIFY: 10",
            "  DMA read: 48  |  DMA write: 51  |  MSI: 12",
            "  TX Forwarded: 4  |  RX Injected: 11",
            "",
            "瓶颈: VCS RTL 仿真速度 (~5000 TLP/min)",
            "  200 个 ping 需要数小时完成",
            "  非代码问题，可通过 VCS 编译优化或换仿真器改善",
        ])

    out = "/home/ubuntu/ryan/software/cosim-platform/docs/cosim_issue_report.pptx"
    prs.save(out)
    print(f"PPT 已生成: {out}")


if __name__ == "__main__":
    main()
