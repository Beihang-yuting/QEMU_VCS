#!/usr/bin/env python3
"""生成 QEMU-VCS 软硬件协同仿真平台设计 PPT"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# 颜色方案
BG_DARK = RGBColor(0x1A, 0x1A, 0x2E)
BG_CARD = RGBColor(0x25, 0x25, 0x3D)
TEXT_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_LIGHT = RGBColor(0xCC, 0xCC, 0xCC)
TEXT_SUBTITLE = RGBColor(0xAA, 0xAA, 0xBB)
ACCENT_BLUE = RGBColor(0x4E, 0x9A, 0xF5)
ACCENT_GREEN = RGBColor(0x4C, 0xAF, 0x50)
ACCENT_ORANGE = RGBColor(0xFF, 0x98, 0x00)
ACCENT_PURPLE = RGBColor(0x9C, 0x27, 0xB0)
ACCENT_RED = RGBColor(0xE5, 0x73, 0x73)
ACCENT_CYAN = RGBColor(0x26, 0xC6, 0xDA)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color=None, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color or BG_CARD
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18, color=TEXT_WHITE,
                 bold=False, alignment=PP_ALIGN.LEFT, font_name="Microsoft YaHei"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=14, color=TEXT_LIGHT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Microsoft YaHei"
        p.space_after = Pt(6)
        p.level = 0
    return txBox


# ============================================================
# Slide 1: 封面
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide, BG_DARK)
add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.5),
             "QEMU-VCS 软硬件协同仿真平台", font_size=44, bold=True, color=TEXT_WHITE,
             alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(3.2), Inches(11), Inches(1),
             "DPU/智能网卡芯片设计验证方案", font_size=28, color=ACCENT_BLUE,
             alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(5.0), Inches(11), Inches(0.5),
             "2026-04-16  |  设计文档", font_size=18, color=TEXT_SUBTITLE,
             alignment=PP_ALIGN.CENTER)

# ============================================================
# Slide 2: 目标与背景
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "项目目标", font_size=36, bold=True, color=TEXT_WHITE)

# 左侧卡片 - 背景
card = add_shape(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5.2), BG_CARD, ACCENT_BLUE)
add_text_box(slide, Inches(1.2), Inches(1.7), Inches(5), Inches(0.5),
             "背景", font_size=22, bold=True, color=ACCENT_BLUE)
add_bullet_list(slide, Inches(1.2), Inches(2.4), Inches(4.8), Inches(4),
                ["• DPU/智能网卡芯片设计验证阶段",
                 "• Host 端驱动需通过 PCIe 访问 DPU RTL",
                 "• 核心接口：PCIe + 以太网",
                 "• 需要端到端功能验证能力"],
                font_size=16)

# 右侧卡片 - 目标
card = add_shape(slide, Inches(6.8), Inches(1.5), Inches(5.5), Inches(5.2), BG_CARD, ACCENT_GREEN)
add_text_box(slide, Inches(7.2), Inches(1.7), Inches(5), Inches(0.5),
             "目标", font_size=22, bold=True, color=ACCENT_GREEN)
add_bullet_list(slide, Inches(7.2), Inches(2.4), Inches(4.8), Inches(4),
                ["• QEMU 模拟 x86 Host，跑 Linux + DPU 网卡驱动",
                 "• VCS 仿真 DPU RTL（PCIe EP、ETH MAC、数据面）",
                 "• 双模式仿真：快速事务级 + 周期精确",
                 "• 双节点互打：两组 QEMU+VCS 通过 ETH 互通",
                 "• 面向验证工程师和驱动开发者"],
                font_size=16)

# ============================================================
# Slide 3: 三种方案概览
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "方案评估 — 三种架构方案", font_size=36, bold=True, color=TEXT_WHITE)

schemes = [
    ("方案 A", "共享内存 + Socket 同步", ACCENT_GREEN,
     ["QEMU 和 VCS 独立进程", "POSIX 共享内存传 TLP 数据", "Unix Socket 做握手同步",
      "✓ 实现门槛最低", "✓ 进程独立，分别调试", "✓ 容易实现精度切换",
      "✗ 同步开销较大", "✗ 需处理同步边界情况"]),
    ("方案 B", "SystemC/TLM 桥接", ACCENT_BLUE,
     ["SystemC 做中间层", "QEMU → TLM-2.0 Socket 发事务", "SystemC 桥转信号级驱动",
      "✓ 行业标准，成熟度高", "✓ TLM-2.0 天然支持精度切换", "✓ 可复用开源 TLM 模型",
      "✗ QEMU 集成 SystemC 复杂", "✗ 需要 co-sim license"]),
    ("方案 C", "DPI-C 直连 + QEMU 库模式", ACCENT_ORANGE,
     ["QEMU 编译为库 (libqemu)", "VCS 通过 DPI-C 同进程调用", "共享地址空间，零拷贝",
      "✓ 性能最优，无 IPC 开销", "✓ 同步最简单", "✓ 事务和中断路径直接",
      "✗ QEMU 改造量大", "✗ 崩溃会拖垮 VCS"]),
]

for i, (name, subtitle, color, items) in enumerate(schemes):
    left = Inches(0.5 + i * 4.2)
    card = add_shape(slide, left, Inches(1.5), Inches(3.8), Inches(5.5), BG_CARD, color)
    add_text_box(slide, left + Inches(0.3), Inches(1.7), Inches(3.2), Inches(0.5),
                 name, font_size=22, bold=True, color=color)
    add_text_box(slide, left + Inches(0.3), Inches(2.2), Inches(3.2), Inches(0.4),
                 subtitle, font_size=14, color=TEXT_SUBTITLE)
    add_bullet_list(slide, left + Inches(0.3), Inches(2.8), Inches(3.2), Inches(4),
                    items, font_size=13, color=TEXT_LIGHT)

# ============================================================
# Slide 4: 选定方案 A 详解
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "选定方案：共享内存 + Socket 同步", font_size=36, bold=True, color=ACCENT_GREEN)
add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.5),
             "团队 QEMU 经验较新，方案 A 门槛最低，能最快跑通端到端，后续可平滑演进",
             font_size=16, color=TEXT_SUBTITLE)

card = add_shape(slide, Inches(0.8), Inches(1.8), Inches(11.5), Inches(5.0), BG_CARD)
add_text_box(slide, Inches(1.2), Inches(2.0), Inches(11), Inches(4.5),
             ("┌──────────────────────────┐        ┌────────────────────────────┐\n"
              "│      QEMU 进程 (Host)     │        │     VCS 仿真进程 (DPU)      │\n"
              "│                          │        │                            │\n"
              "│  Guest OS (Linux)        │        │   DPU RTL Top              │\n"
              "│    └─ DPU NIC Driver     │        │   ├─ PCIe EP (RTL)         │\n"
              "│         │ MMIO/DMA       │        │   ├─ ETH MAC (RTL)         │\n"
              "│  PCIe RC Device Model    │        │   └─ 数据面逻辑 (RTL)       │\n"
              "│  (自定义 QEMU 设备)       │        │                            │\n"
              "│         │                │        │        │                   │\n"
              "│  Bridge Adapter (C lib)  │        │   Bridge Adapter (DPI-C)   │\n"
              "└─────────┼────────────────┘        └────────┼───────────────────┘\n"
              "          │                                  │\n"
              "   ┌──────┴──────────────────────────────────┴──────┐\n"
              "   │  共享内存 (SHM)        │  同步 Socket (Unix)    │\n"
              "   │  • TLP 请求/响应队列   │  • 事务握手/中断通知   │\n"
              "   │  • DMA 数据缓冲        │  • 时钟同步/模式切换   │\n"
              "   └───────────────────────────────────────────────┘"),
             font_size=12, color=ACCENT_CYAN, font_name="Consolas")

# ============================================================
# Slide 5: 数据流与同步
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "数据流与同步机制", font_size=36, bold=True, color=TEXT_WHITE)

# PCIe MMIO Read 流程
card = add_shape(slide, Inches(0.5), Inches(1.5), Inches(12.3), Inches(2.2), BG_CARD, ACCENT_BLUE)
add_text_box(slide, Inches(0.9), Inches(1.6), Inches(5), Inches(0.4),
             "PCIe MMIO 读事务流程", font_size=18, bold=True, color=ACCENT_BLUE)
add_text_box(slide, Inches(0.9), Inches(2.1), Inches(11.5), Inches(1.5),
             ("① Guest 驱动 MMIO Read  →  ② QEMU 拦截, 暂停 vCPU  →  ③ 构造 MRd TLP, 写入 SHM\n"
              "→  ④ Socket 通知 VCS  →  ⑤ DPI-C 读取 TLP, 驱动 RTL  →  ⑥ RTL 返回数据, 写入 SHM\n"
              "→  ⑦ Socket 通知 QEMU  →  ⑧ 恢复 vCPU, 返回数据给 Guest"),
             font_size=14, color=TEXT_LIGHT)

# 双模式
card1 = add_shape(slide, Inches(0.5), Inches(4.0), Inches(5.8), Inches(3.0), BG_CARD, ACCENT_GREEN)
add_text_box(slide, Inches(0.9), Inches(4.1), Inches(5), Inches(0.4),
             "快速模式（事务级）", font_size=18, bold=True, color=ACCENT_GREEN)
add_bullet_list(slide, Inches(0.9), Inches(4.6), Inches(5), Inches(2.2),
                ["• 同步粒度：每个 PCIe 事务",
                 "• QEMU 自由运行，遇 PCIe 访问暂停",
                 "• VCS 快进到事务完成",
                 "• 适合：驱动加载、配置、功能验证",
                 "• 速度：~1000-10000 事务/秒"],
                font_size=13)

card2 = add_shape(slide, Inches(7.0), Inches(4.0), Inches(5.8), Inches(3.0), BG_CARD, ACCENT_ORANGE)
add_text_box(slide, Inches(7.4), Inches(4.1), Inches(5), Inches(0.4),
             "精确模式（周期级）", font_size=18, bold=True, color=ACCENT_ORANGE)
add_bullet_list(slide, Inches(7.4), Inches(4.6), Inches(5), Inches(2.2),
                ["• 同步粒度：每 N 个时钟周期",
                 "• VCS 与 QEMU 锁步推进",
                 "• 保证时序一致性",
                 "• 适合：时序调试、波形分析",
                 "• 速度：~10-100 事务/秒"],
                font_size=13)

# ============================================================
# Slide 6: 共享内存数据结构
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "共享内存数据结构与 API", font_size=36, bold=True, color=TEXT_WHITE)

# SHM 布局
card = add_shape(slide, Inches(0.5), Inches(1.5), Inches(6), Inches(5.5), BG_CARD, ACCENT_BLUE)
add_text_box(slide, Inches(0.9), Inches(1.6), Inches(5), Inches(0.4),
             "PCIe SHM 布局 (64MB)", font_size=18, bold=True, color=ACCENT_BLUE)
add_text_box(slide, Inches(0.9), Inches(2.2), Inches(5.5), Inches(4.5),
             ("┌─────────────────────────────┐\n"
              "│ [0x0000] 控制区 (4KB)        │\n"
              "│  magic / version / mode      │\n"
              "│  就绪标志 / 仿真时间 / IRQ   │\n"
              "├─────────────────────────────┤\n"
              "│ [0x1000] 请求队列 (256KB)    │\n"
              "│  Ring Buffer: Host→Device    │\n"
              "│  Entry: type/addr/data/tag   │\n"
              "├─────────────────────────────┤\n"
              "│ [0x41000] 响应队列 (256KB)   │\n"
              "│  Ring Buffer: Device→Host    │\n"
              "│  Entry: tag/status/data      │\n"
              "├─────────────────────────────┤\n"
              "│ [0x81000] DMA 数据区 (63MB)  │\n"
              "│  按 4KB 页对齐，多 DMA 通道   │\n"
              "└─────────────────────────────┘"),
             font_size=12, color=ACCENT_CYAN, font_name="Consolas")

# API
card = add_shape(slide, Inches(6.8), Inches(1.5), Inches(6), Inches(2.5), BG_CARD, ACCENT_GREEN)
add_text_box(slide, Inches(7.2), Inches(1.6), Inches(5), Inches(0.4),
             "QEMU 侧 Bridge API (C)", font_size=16, bold=True, color=ACCENT_GREEN)
add_text_box(slide, Inches(7.2), Inches(2.1), Inches(5.5), Inches(1.8),
             ("bridge_init() / bridge_connect()\n"
              "bridge_send_tlp() / bridge_wait_completion()\n"
              "bridge_poll_dma() / bridge_complete_dma()\n"
              "bridge_inject_msi() / bridge_poll_irq()\n"
              "bridge_set_mode() / bridge_advance_time()"),
             font_size=12, color=TEXT_LIGHT, font_name="Consolas")

card = add_shape(slide, Inches(6.8), Inches(4.3), Inches(6), Inches(2.7), BG_CARD, ACCENT_ORANGE)
add_text_box(slide, Inches(7.2), Inches(4.4), Inches(5), Inches(0.4),
             "VCS 侧 Bridge API (DPI-C)", font_size=16, bold=True, color=ACCENT_ORANGE)
add_text_box(slide, Inches(7.2), Inches(4.9), Inches(5.5), Inches(2),
             ("bridge_vcs_init(shm_name)\n"
              "bridge_vcs_poll_tlp(type, addr, data, len)\n"
              "bridge_vcs_send_completion(tag, data, len)\n"
              "bridge_vcs_trigger_dma(host_addr, data, len)\n"
              "bridge_vcs_raise_msi(vec)\n"
              "bridge_vcs_sync_step(clks)"),
             font_size=12, color=TEXT_LIGHT, font_name="Consolas")

# ============================================================
# Slide 7: 双节点互打架构
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "双节点互打架构", font_size=36, bold=True, color=TEXT_WHITE)
add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.5),
             "两组 QEMU+VCS 通过 ETH 共享内存实现端到端数据互通",
             font_size=16, color=TEXT_SUBTITLE)

card = add_shape(slide, Inches(0.5), Inches(1.8), Inches(12.3), Inches(3.2), BG_CARD)
add_text_box(slide, Inches(0.9), Inches(1.9), Inches(11.5), Inches(3),
             ("                    节点 A                                               节点 B\n"
              "  ┌───────────────────────────────┐                  ┌───────────────────────────────┐\n"
              "  │ QEMU-A    ↔  [PCIe SHM-A]  ↔  VCS-A │    ↔    │ VCS-B  ↔  [PCIe SHM-B]  ↔  QEMU-B │\n"
              "  │ (Host)       (TLP 通信)      (DPU)  │          │ (DPU)     (TLP 通信)      (Host)  │\n"
              "  └───────────────────────────────┘                  └───────────────────────────────┘\n"
              "                                      │              │\n"
              "                                ┌─────┴──────────────┴─────┐\n"
              "                                │      ETH SHM (新增)       │\n"
              "                                │  A→B 帧队列 / B→A 帧队列  │\n"
              "                                │  链路速率 / 延迟 / 丢包    │\n"
              "                                └──────────────────────────┘"),
             font_size=12, color=ACCENT_CYAN, font_name="Consolas")

# 三块 SHM 说明
items_data = [
    ("PCIe SHM-A", "QEMU-A ↔ VCS-A\nPCIe TLP 事务通信", ACCENT_BLUE),
    ("ETH SHM", "VCS-A ↔ VCS-B\n以太网帧传输", ACCENT_ORANGE),
    ("PCIe SHM-B", "QEMU-B ↔ VCS-B\nPCIe TLP 事务通信", ACCENT_BLUE),
]
for i, (title, desc, color) in enumerate(items_data):
    left = Inches(0.5 + i * 4.3)
    c = add_shape(slide, left, Inches(5.3), Inches(3.8), Inches(1.8), BG_CARD, color)
    add_text_box(slide, left + Inches(0.3), Inches(5.4), Inches(3.2), Inches(0.4),
                 title, font_size=18, bold=True, color=color)
    add_text_box(slide, left + Inches(0.3), Inches(5.9), Inches(3.2), Inches(1),
                 desc, font_size=14, color=TEXT_LIGHT)

# ============================================================
# Slide 8: QEMU 自定义设备集成
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "QEMU 自定义 PCIe 设备集成", font_size=36, bold=True, color=TEXT_WHITE)

# 左侧 - 源码位置
card = add_shape(slide, Inches(0.5), Inches(1.5), Inches(6), Inches(2.5), BG_CARD, ACCENT_BLUE)
add_text_box(slide, Inches(0.9), Inches(1.6), Inches(5), Inches(0.4),
             "In-tree 编译集成", font_size=18, bold=True, color=ACCENT_BLUE)
add_text_box(slide, Inches(0.9), Inches(2.2), Inches(5.5), Inches(1.5),
             ("qemu-source/\n"
              "├── hw/net/cosim_pcie_rc.c       ← 设备实现\n"
              "├── include/hw/net/cosim_pcie_rc.h\n"
              "└── hw/net/meson.build            ← 添加编译目标\n\n"
              "启动: -device cosim-pcie-rc,shm_name=...,sock_path=..."),
             font_size=12, color=TEXT_LIGHT, font_name="Consolas")

# 右侧 - 核心逻辑
card = add_shape(slide, Inches(6.8), Inches(1.5), Inches(6), Inches(2.5), BG_CARD, ACCENT_GREEN)
add_text_box(slide, Inches(7.2), Inches(1.6), Inches(5), Inches(0.4),
             "设备模型核心逻辑", font_size=18, bold=True, color=ACCENT_GREEN)
add_bullet_list(slide, Inches(7.2), Inches(2.2), Inches(5.3), Inches(1.5),
                ["• MMIO 读：拦截 → MRd TLP → Bridge → 等待 → 返回",
                 "• MMIO 写：拦截 → MWr TLP → Bridge → 发送",
                 "• 初始化：注册 BAR + PCI ID + Bridge 连接",
                 "• MSI 中断：轮询检测 → msi_notify() 注入"],
                font_size=13)

# 下方 - 一键脚本
card = add_shape(slide, Inches(0.5), Inches(4.3), Inches(12.3), Inches(2.8), BG_CARD, ACCENT_ORANGE)
add_text_box(slide, Inches(0.9), Inches(4.4), Inches(5), Inches(0.4),
             "一键搭建脚本 setup_cosim_qemu.sh", font_size=18, bold=True, color=ACCENT_ORANGE)
add_text_box(slide, Inches(0.9), Inches(5.0), Inches(11.5), Inches(2),
             ("① 安装依赖 (build-essential, ninja, meson, libglib2.0-dev ...)\n"
              "② 下载 QEMU 源码 (git clone --branch v9.2.0 --depth 1)\n"
              "③ 注入自定义设备代码 (cp + patch meson.build)\n"
              "④ 编译 Bridge 库 (cmake → libcosim_bridge.so)\n"
              "⑤ 编译 QEMU (./configure --extra-cflags/ldflags → ninja)\n"
              "⑥ 准备 Guest 镜像 (预装 DPU 网卡驱动)\n"
              "⑦ 生成启动脚本 run_cosim.sh"),
             font_size=13, color=TEXT_LIGHT, font_name="Consolas")

# ============================================================
# Slide 9: 调试能力
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "调试与工具链", font_size=36, bold=True, color=TEXT_WHITE)

# 验证工程师
card = add_shape(slide, Inches(0.5), Inches(1.4), Inches(5.8), Inches(2.8), BG_CARD, ACCENT_BLUE)
add_text_box(slide, Inches(0.9), Inches(1.5), Inches(5), Inches(0.4),
             "验证工程师", font_size=20, bold=True, color=ACCENT_BLUE)
add_bullet_list(slide, Inches(0.9), Inches(2.1), Inches(5), Inches(2),
                ["• VCS 原生波形 (VPD/FSDB)",
                 "• 事务级追踪日志 (CSV/JSON)",
                 "• SVA 断言 + 覆盖率收集",
                 "• Bridge 事务也记录到波形"],
                font_size=14)

# 驱动开发者
card = add_shape(slide, Inches(7.0), Inches(1.4), Inches(5.8), Inches(2.8), BG_CARD, ACCENT_GREEN)
add_text_box(slide, Inches(7.4), Inches(1.5), Inches(5), Inches(0.4),
             "驱动开发者", font_size=20, bold=True, color=ACCENT_GREEN)
add_bullet_list(slide, Inches(7.4), Inches(2.1), Inches(5), Inches(2),
                ["• GDB 调试 Guest 内核/驱动",
                 "• QEMU Monitor 查看 PCIe/BAR",
                 "• Guest dmesg 实时输出",
                 "• Bridge 通信延迟统计"],
                font_size=14)

# 统一控制台
card = add_shape(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.6), BG_CARD, ACCENT_PURPLE)
add_text_box(slide, Inches(0.9), Inches(4.6), Inches(5), Inches(0.4),
             "统一调试控制台 cosim_cli.py", font_size=20, bold=True, color=ACCENT_PURPLE)
add_text_box(slide, Inches(0.9), Inches(5.2), Inches(11.5), Inches(1.8),
             ("$ cosim start --qemu-image linux.qcow2 --vcs-top dpu_top      # 启动仿真\n"
              "$ cosim mode precise --trigger \"addr==0xFE00_0000\"            # 条件触发切精确模式\n"
              "$ cosim trace dump --format csv --filter \"type==MRd\"          # 事务追踪\n"
              "$ cosim stats --last 1000                                     # 延迟分布统计"),
             font_size=13, color=TEXT_LIGHT, font_name="Consolas")

# ============================================================
# Slide 10: 项目结构
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "项目结构", font_size=36, bold=True, color=TEXT_WHITE)

card = add_shape(slide, Inches(0.5), Inches(1.3), Inches(12.3), Inches(5.8), BG_CARD)
add_text_box(slide, Inches(0.9), Inches(1.4), Inches(11.5), Inches(5.5),
             ("cosim-platform/\n"
              "├── bridge/                          # 通信层核心库\n"
              "│   ├── common/                      #   SHM 布局、环形缓冲、协议定义\n"
              "│   ├── qemu/                        #   QEMU 侧 Bridge Adapter (C)\n"
              "│   ├── vcs/                         #   VCS 侧 Bridge Adapter (DPI-C)\n"
              "│   └── eth/                         #   ETH 共享内存（双节点互打）\n"
              "│\n"
              "├── qemu-plugin/                     # QEMU 自定义 PCIe RC 设备模型\n"
              "│   ├── cosim_pcie_rc.c / .h\n"
              "│   └── Kconfig\n"
              "│\n"
              "├── vcs-tb/                          # VCS Testbench\n"
              "│   ├── tb_top.sv / pcie_ep_driver.sv / eth_mac_driver.sv\n"
              "│   └── uvm/                         #   UVM 环境（可选）\n"
              "│\n"
              "├── tools/                           # cosim_cli.py / trace_analyzer.py / launch.py\n"
              "├── tests/                           # 测试用例 (BAR/DMA/MSI/双节点 ping)\n"
              "├── setup_cosim_qemu.sh              # 一键搭建\n"
              "├── run_cosim.sh                     # 启动脚本\n"
              "└── CMakeLists.txt + Makefile"),
             font_size=13, color=TEXT_LIGHT, font_name="Consolas")

# ============================================================
# Slide 11: 分阶段交付
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "分阶段交付计划", font_size=36, bold=True, color=TEXT_WHITE)

phases = [
    ("P1", "单节点 PCIe\n通路打通", ACCENT_GREEN,
     ["bridge/common 实现", "QEMU PCIe RC 设备", "VCS 最小 testbench", "快速模式同步"],
     "验收：BAR 读写通过"),
    ("P2", "DMA + 中断\n+ 精确模式", ACCENT_BLUE,
     ["DMA 数据通路", "MSI/MSI-X 中断", "周期级锁步", "模式切换"],
     "验收：DMA+中断通过"),
    ("P3", "双节点\nETH 互打", ACCENT_ORANGE,
     ["ETH SHM 帧队列", "ETH MAC 驱动集成", "双 VCS 时间同步", "双节点启动编排"],
     "验收：ping 端到端通过"),
    ("P4", "调试工具\n与生产化", ACCENT_PURPLE,
     ["cosim_cli 控制台", "trace_analyzer", "CI 回归流水线", "链路模型增强"],
     "验收：团队可独立使用"),
]

for i, (label, title, color, items, acceptance) in enumerate(phases):
    left = Inches(0.3 + i * 3.25)
    card = add_shape(slide, left, Inches(1.4), Inches(3.0), Inches(5.5), BG_CARD, color)

    # Phase label
    lbl = add_shape(slide, left + Inches(0.2), Inches(1.6), Inches(0.7), Inches(0.5), color)
    add_text_box(slide, left + Inches(0.2), Inches(1.6), Inches(0.7), Inches(0.5),
                 label, font_size=16, bold=True, color=TEXT_WHITE, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, left + Inches(1.0), Inches(1.6), Inches(1.8), Inches(0.7),
                 title, font_size=14, bold=True, color=color)

    add_bullet_list(slide, left + Inches(0.3), Inches(2.6), Inches(2.5), Inches(2.8),
                    [f"• {item}" for item in items], font_size=12, color=TEXT_LIGHT)

    # Acceptance
    add_shape(slide, left + Inches(0.2), Inches(5.5), Inches(2.6), Inches(1.1),
              RGBColor(0x1E, 0x3A, 0x1E), color)
    add_text_box(slide, left + Inches(0.3), Inches(5.7), Inches(2.4), Inches(0.8),
                 acceptance, font_size=11, bold=True, color=color, alignment=PP_ALIGN.CENTER)

# ============================================================
# Slide 12: 总结
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_text_box(slide, Inches(1), Inches(1.0), Inches(11), Inches(1.2),
             "总结", font_size=44, bold=True, color=TEXT_WHITE, alignment=PP_ALIGN.CENTER)

summary_items = [
    ("方案", "共享内存 + Socket 同步（门槛低、可演进）", ACCENT_GREEN),
    ("架构", "QEMU(Host) ↔ SHM+Socket ↔ VCS(DPU RTL)", ACCENT_BLUE),
    ("模式", "事务级快速 + 周期精确，可运行时切换", ACCENT_ORANGE),
    ("扩展", "双节点 ETH SHM 互打，端到端验证", ACCENT_PURPLE),
    ("工具", "一键搭建 + 统一调试控制台 + 事务追踪", ACCENT_CYAN),
]

for i, (label, desc, color) in enumerate(summary_items):
    top = Inches(2.5 + i * 0.9)
    lbl_shape = add_shape(slide, Inches(2), top, Inches(1.5), Inches(0.6), color)
    add_text_box(slide, Inches(2), top + Inches(0.05), Inches(1.5), Inches(0.5),
                 label, font_size=16, bold=True, color=TEXT_WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(3.8), top + Inches(0.05), Inches(7), Inches(0.5),
                 desc, font_size=16, color=TEXT_LIGHT)

# Save
output_path = "/home/ubuntu/ryan/software/QEMU-VCS-CoSim-Platform-Design.pptx"
prs.save(output_path)
print(f"PPT saved to: {output_path}")
