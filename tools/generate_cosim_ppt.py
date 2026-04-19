#!/usr/bin/env python3
"""Generate QEMU-VCS CoSim Platform comprehensive PPT report in Chinese."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


# ── Color scheme ──────────────────────────────────────────────────────────
DARK_BLUE = RGBColor(0x00, 0x2B, 0x5C)
MEDIUM_BLUE = RGBColor(0x00, 0x56, 0xA0)
LIGHT_BLUE = RGBColor(0x00, 0x7B, 0xC0)
ACCENT_BLUE = RGBColor(0xD6, 0xE8, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MEDIUM_GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
GREEN = RGBColor(0x00, 0x80, 0x00)
RED = RGBColor(0xCC, 0x00, 0x00)
ORANGE = RGBColor(0xFF, 0x8C, 0x00)
HEADER_BG = RGBColor(0x00, 0x3D, 0x7A)
ROW_ALT = RGBColor(0xE8, 0xF0, 0xFA)

CN_FONT = "SimHei"
EN_FONT = "Arial"
CODE_FONT = "Courier New"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ── Helper functions ──────────────────────────────────────────────────────

def add_background(slide, color=WHITE):
    """Set slide background color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_bar(slide, title_text, subtitle_text=None):
    """Add a dark blue title bar at the top of the slide."""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(1.2)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK_BLUE
    bar.line.fill.background()

    tf = bar.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = CN_FONT
    p.alignment = PP_ALIGN.LEFT
    tf.margin_left = Inches(0.5)
    tf.margin_top = Inches(0.15)

    if subtitle_text:
        p2 = tf.add_paragraph()
        p2.text = subtitle_text
        p2.font.size = Pt(16)
        p2.font.color.rgb = ACCENT_BLUE
        p2.font.name = CN_FONT
        p2.alignment = PP_ALIGN.LEFT

    # Bottom accent line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.2), SLIDE_W, Inches(0.05)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = LIGHT_BLUE
    line.line.fill.background()


def add_footer(slide, page_num, total=25):
    """Add footer with page number and project name."""
    tf = slide.shapes.add_textbox(
        Inches(0.5), Inches(7.0), Inches(5), Inches(0.4)
    ).text_frame
    p = tf.paragraphs[0]
    p.text = "QEMU-VCS CoSim Platform"
    p.font.size = Pt(10)
    p.font.color.rgb = MEDIUM_GRAY
    p.font.name = EN_FONT

    tf2 = slide.shapes.add_textbox(
        Inches(11), Inches(7.0), Inches(2), Inches(0.4)
    ).text_frame
    p2 = tf2.paragraphs[0]
    p2.text = f"{page_num} / {total}"
    p2.font.size = Pt(10)
    p2.font.color.rgb = MEDIUM_GRAY
    p2.font.name = EN_FONT
    p2.alignment = PP_ALIGN.RIGHT


def add_content_box(slide, left, top, width, height):
    """Add a text box and return the text frame."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    return tf


def add_bullet(tf, text, level=0, font_size=18, bold=False, color=DARK_GRAY,
               font_name=CN_FONT):
    """Add a bullet point paragraph to a text frame."""
    if len(tf.paragraphs) == 1 and tf.paragraphs[0].text == "":
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = text
    p.level = level
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.space_after = Pt(4)
    p.space_before = Pt(2)
    return p


def add_code_block(slide, left, top, width, height, code_text):
    """Add a code block (gray background box with monospace text)."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
    shape.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    shape.line.width = Pt(1)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.1)
    for i, line in enumerate(code_text.strip().split("\n")):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(12)
        p.font.name = CODE_FONT
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(1)
    return shape


def add_diagram_box(slide, left, top, width, height, text,
                    bg_color=ACCENT_BLUE, border_color=MEDIUM_BLUE,
                    font_size=14):
    """Add a styled box for diagrams."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(2)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.font.name = CN_FONT
    p.alignment = PP_ALIGN.CENTER
    return shape


def add_arrow(slide, left, top, width, height=None):
    """Add a right arrow shape."""
    if height is None:
        height = Inches(0.35)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = MEDIUM_BLUE
    shape.line.fill.background()
    return shape


def add_down_arrow(slide, left, top, width=Inches(0.35), height=Inches(0.5)):
    """Add a down arrow shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.DOWN_ARROW, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = MEDIUM_BLUE
    shape.line.fill.background()
    return shape


def add_table(slide, left, top, width, height, rows, cols, data,
              col_widths=None):
    """Add a table with header styling."""
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = data[r][c]
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(13)
                paragraph.font.name = CN_FONT
                paragraph.alignment = PP_ALIGN.CENTER

                if r == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
                else:
                    paragraph.font.color.rgb = DARK_GRAY

            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = HEADER_BG
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = ROW_ALT
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE

    return table_shape


def add_section_divider(slide, section_num, section_title, subtitle=""):
    """Create a section divider slide."""
    add_background(slide, DARK_BLUE)

    # Section number circle
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(5.7), Inches(2.0), Inches(2.0), Inches(2.0)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = LIGHT_BLUE
    circle.line.fill.background()
    tf = circle.text_frame
    p = tf.paragraphs[0]
    p.text = str(section_num)
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Section title
    tf2 = add_content_box(slide, Inches(1), Inches(4.3), Inches(11), Inches(1.5))
    p2 = tf2.paragraphs[0]
    p2.text = section_title
    p2.font.size = Pt(40)
    p2.font.bold = True
    p2.font.color.rgb = WHITE
    p2.font.name = CN_FONT
    p2.alignment = PP_ALIGN.CENTER

    if subtitle:
        p3 = tf2.add_paragraph()
        p3.text = subtitle
        p3.font.size = Pt(20)
        p3.font.color.rgb = ACCENT_BLUE
        p3.font.name = CN_FONT
        p3.alignment = PP_ALIGN.CENTER


TOTAL_PAGES = 42
page = 0


def new_slide():
    global page
    page += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    add_background(slide)
    return slide, page


# ======================================================================
# SLIDE 1 -- Cover
# ======================================================================
slide, pg = new_slide()
add_background(slide, DARK_BLUE)

tf = add_content_box(slide, Inches(1), Inches(1.5), Inches(11.3), Inches(2.5))
p = tf.paragraphs[0]
p.text = "CoSim \u5e73\u53f0"
p.font.size = Pt(54)
p.font.bold = True
p.font.color.rgb = WHITE
p.font.name = CN_FONT
p.alignment = PP_ALIGN.CENTER

p2 = tf.add_paragraph()
p2.text = "QEMU-VCS PCIe \u534f\u540c\u4eff\u771f"
p2.font.size = Pt(36)
p2.font.color.rgb = ACCENT_BLUE
p2.font.name = CN_FONT
p2.alignment = PP_ALIGN.CENTER

line = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(4), Inches(4.2), Inches(5.3), Inches(0.04)
)
line.fill.solid()
line.fill.fore_color.rgb = LIGHT_BLUE
line.line.fill.background()

tf2 = add_content_box(slide, Inches(1), Inches(4.5), Inches(11.3), Inches(2))
for text in ["PCIe Root Complex / Endpoint \u534f\u540c\u4eff\u771f\u5e73\u53f0",
             "\u65e5\u671f: 2026-04-18"]:
    add_bullet(tf2, text, font_size=22, color=WHITE, bold=False)
    tf2.paragraphs[-1].alignment = PP_ALIGN.CENTER

# ======================================================================
# SLIDE 2 -- System Architecture Overview
# ======================================================================
slide, pg = new_slide()
add_title_bar(slide, "\u7cfb\u7edf\u67b6\u6784\u603b\u89c8")
add_footer(slide, pg, TOTAL_PAGES)

add_diagram_box(slide, Inches(0.5), Inches(1.8), Inches(3.5), Inches(2.5),
                "QEMU\n(Guest OS)\n\ncosim-pcie-rc \u8bbe\u5907\nMMIO / Config \u56de\u8c03",
                bg_color=RGBColor(0xE3, 0xF2, 0xFD), font_size=16)

add_arrow(slide, Inches(4.1), Inches(2.8), Inches(1.2))

add_diagram_box(slide, Inches(5.4), Inches(1.8), Inches(3.0), Inches(2.5),
                "Bridge \u5c42\n\nSHM Ring Buffer\nUnix Socket \u63a7\u5236\n\u547d\u4ee4\u901a\u9053 + TLP \u901a\u9053",
                bg_color=RGBColor(0xFF, 0xF8, 0xE1), border_color=ORANGE,
                font_size=16)

add_arrow(slide, Inches(8.5), Inches(2.8), Inches(1.2))

add_diagram_box(slide, Inches(9.8), Inches(1.8), Inches(3.0), Inches(2.5),
                "VCS\n(EP Stub)\n\npcie_ep_stub.sv\nDPI-C \u63a5\u53e3\nVirtio \u5bc4\u5b58\u5668",
                bg_color=RGBColor(0xE8, 0xF5, 0xE9), border_color=GREEN,
                font_size=16)

tf = add_content_box(slide, Inches(0.5), Inches(4.8), Inches(12), Inches(2.2))
add_bullet(tf, "\u4e09\u5927\u6838\u5fc3\u7ec4\u4ef6:", font_size=18, bold=True,
           color=DARK_BLUE)
add_bullet(tf,
           "QEMU cosim-pcie-rc: PCIe Root Complex \u6a21\u62df\u8bbe\u5907, \u62e6\u622a Guest MMIO/Config \u8bbf\u95ee",
           level=1, font_size=16)
add_bullet(tf,
           "Bridge \u5c42: \u5171\u4eab\u5185\u5b58 Ring Buffer + Unix Socket \u5b9e\u73b0\u8de8\u8fdb\u7a0b TLP \u4f20\u8f93",
           level=1, font_size=16)
add_bullet(tf,
           "VCS pcie_ep_stub: SystemVerilog Endpoint \u6a21\u578b, \u901a\u8fc7 DPI-C \u4e0e Bridge \u901a\u4fe1",
           level=1, font_size=16)

# ======================================================================
# SLIDE 3 -- Data Path: TLP Flow
# ======================================================================
slide, pg = new_slide()
add_title_bar(slide, "\u6570\u636e\u901a\u8def: TLP \u6d41\u5411\u56fe")
add_footer(slide, pg, TOTAL_PAGES)

tf_label = add_content_box(slide, Inches(0.3), Inches(1.5), Inches(3), Inches(0.4))
add_bullet(tf_label, "MMIO \u901a\u8def (MRd/MWr):", font_size=16, bold=True,
           color=DARK_BLUE)

boxes_top = [
    ("Guest\nMMIO", Inches(0.3), Inches(2.0)),
    ("QEMU\nmmio_ops", Inches(2.6), Inches(2.0)),
    ("bridge_\nsend_tlp", Inches(4.9), Inches(2.0)),
    ("SHM\nRing", Inches(7.2), Inches(2.0)),
    ("VCS DPI\ncosim_recv", Inches(9.5), Inches(2.0)),
    ("EP Stub\n\u5bc4\u5b58\u5668", Inches(11.5), Inches(2.0)),
]
for txt, l, t in boxes_top:
    add_diagram_box(slide, l, t, Inches(1.8), Inches(0.9), txt, font_size=12)

for i in range(len(boxes_top) - 1):
    l = boxes_top[i][1] + Inches(1.85)
    add_arrow(slide, l, Inches(2.25), Inches(0.65), Inches(0.25))

tf_label2 = add_content_box(slide, Inches(0.3), Inches(3.4), Inches(3),
                            Inches(0.4))
add_bullet(tf_label2, "DMA \u53cd\u5411\u901a\u8def:", font_size=16, bold=True,
           color=DARK_BLUE)

boxes_bot = [
    ("VCS EP\nDMA Req", Inches(0.3), Inches(3.9)),
    ("bridge_\nsend_dma", Inches(2.6), Inches(3.9)),
    ("SHM\nRing", Inches(4.9), Inches(3.9)),
    ("QEMU\nDMA \u5904\u7406", Inches(7.2), Inches(3.9)),
    ("address_\nspace_write", Inches(9.5), Inches(3.9)),
    ("Guest\nMemory", Inches(11.5), Inches(3.9)),
]
for txt, l, t in boxes_bot:
    add_diagram_box(slide, l, t, Inches(1.8), Inches(0.9), txt,
                    bg_color=RGBColor(0xE8, 0xF5, 0xE9), border_color=GREEN,
                    font_size=12)

for i in range(len(boxes_bot) - 1):
    l = boxes_bot[i][1] + Inches(1.85)
    add_arrow(slide, l, Inches(4.15), Inches(0.65), Inches(0.25))

tf3 = add_content_box(slide, Inches(0.5), Inches(5.3), Inches(12), Inches(1.5))
add_bullet(tf3,
           "TLP \u683c\u5f0f: \u6807\u51c6 PCIe TLP Header (3DW/4DW) + Data Payload",
           font_size=15)
add_bullet(tf3,
           "Ring Buffer: \u65e0\u9501\u5355\u751f\u4ea7\u8005/\u5355\u6d88\u8d39\u8005, head/tail \u539f\u5b50\u64cd\u4f5c",
           font_size=15)
add_bullet(tf3,
           "Completion: MRd \u8bf7\u6c42\u7531 VCS \u8fd4\u56de CplD, \u901a\u8fc7\u76f8\u540c Ring Buffer \u56de\u4f20",
           font_size=15)

# ======================================================================
# SLIDE 4 -- BAR0 Memory Map
# ======================================================================
slide, pg = new_slide()
add_title_bar(slide, "BAR0 \u5185\u5b58\u6620\u5c04")
add_footer(slide, pg, TOTAL_PAGES)

data = [
    ["\u5730\u5740\u8303\u56f4", "\u5927\u5c0f", "\u529f\u80fd", "\u8bf4\u660e"],
    ["0x000 - 0x1FF", "512B", "\u901a\u7528\u6d4b\u8bd5\u5bc4\u5b58\u5668",
     "devmem_test \u8bfb\u5199\u9a8c\u8bc1\u533a\u57df"],
    ["0x1000 - 0x1FFF", "4KB", "virtio common_cfg",
     "Virtio PCI \u901a\u7528\u914d\u7f6e"],
    ["0x2000 - 0x2FFF", "4KB", "virtio notify",
     "\u961f\u5217\u901a\u77e5\u5bc4\u5b58\u5668"],
    ["0x3000 - 0x3FFF", "4KB", "virtio ISR",
     "\u4e2d\u65ad\u72b6\u6001\u5bc4\u5b58\u5668"],
    ["0x4000 - 0x4FFF", "4KB", "virtio device_cfg",
     "\u8bbe\u5907\u7279\u5b9a\u914d\u7f6e (MAC\u7b49)"],
]
add_table(slide, Inches(1), Inches(1.8), Inches(11), Inches(3.0), 6, 4, data,
          col_widths=[Inches(2.5), Inches(1.5), Inches(3), Inches(4)])

tf = add_content_box(slide, Inches(0.5), Inches(5.2), Inches(12), Inches(1.5))
add_bullet(tf,
           "BAR0 \u603b\u5927\u5c0f: 32KB (0x8000), \u7531 BAR0_SIZE_MASK \u5b9a\u4e49",
           font_size=16, bold=True)
add_bullet(tf,
           "BAR Sizing \u673a\u5236: Guest \u5199 0xFFFFFFFF \u2192 \u8bfb\u56de ~(SIZE-1) & 0xFFFFFFF0",
           font_size=15)
add_bullet(tf,
           "\u6240\u6709\u5bc4\u5b58\u5668\u8bbf\u95ee\u5747\u4e3a Little-Endian MMIO \u8bed\u4e49",
           font_size=15)

# ======================================================================
# SLIDE 5 -- Key File List
# ======================================================================
slide, pg = new_slide()
add_title_bar(slide, "\u5173\u952e\u6587\u4ef6\u6e05\u5355")
add_footer(slide, pg, TOTAL_PAGES)

data = [
    ["\u8def\u5f84", "\u8bed\u8a00", "\u8bf4\u660e"],
    ["qemu-plugin/cosim_pcie_rc.c", "C",
     "QEMU PCIe RC \u8bbe\u5907\u6a21\u578b (MMIO/Config/DMA)"],
    ["bridge/cosim_bridge.c", "C",
     "Bridge \u6838\u5fc3: Ring Buffer + SHM \u7ba1\u7406"],
    ["bridge/cosim_bridge.h", "C",
     "Bridge \u5934\u6587\u4ef6: TLP \u7ed3\u6784\u4f53\u5b9a\u4e49"],
    ["vcs-tb/pcie_ep_stub.sv", "SV",
     "VCS Endpoint Stub (\u5bc4\u5b58\u5668 + Virtio)"],
    ["vcs-tb/cosim_dpi.c", "C",
     "DPI-C \u63a5\u53e3: SV \u2194 Bridge \u6865\u63a5"],
    ["vcs-tb/cosim_dpi.sv", "SV", "DPI-C \u58f0\u660e\u4e0e import"],
    ["scripts/run_cosim.sh", "Bash",
     "\u534f\u540c\u4eff\u771f\u542f\u52a8\u811a\u672c"],
    ["scripts/devmem_test.sh", "Bash",
     "MMIO \u8bfb\u5199\u6d4b\u8bd5\u811a\u672c"],
    ["scripts/virtio_reg_test.sh", "Bash",
     "Virtio \u5bc4\u5b58\u5668\u9a8c\u8bc1\u811a\u672c"],
]
add_table(slide, Inches(0.5), Inches(1.8), Inches(12.3), Inches(4.5), 10, 3,
          data, col_widths=[Inches(4.5), Inches(1.3), Inches(6.5)])

# ======================================================================
# SECTION DIVIDER -- Part II
# ======================================================================
slide, pg = new_slide()
add_section_divider(slide, "II",
                    "\u7b2c\u4e8c\u90e8\u5206: \u5b9e\u73b0\u9636\u6bb5",
                    "Phase 0 \u2192 Phase 2b \u9010\u6b65\u63a8\u8fdb")
add_footer(slide, pg, TOTAL_PAGES)

# ======================================================================
# SLIDE 7 -- Phase 0: Basic MMIO Path
# ======================================================================
slide, pg = new_slide()
add_title_bar(slide,
              "Phase 0: \u57fa\u7840 MMIO \u901a\u8def",
              "MRd / MWr TLP \u8f6c\u53d1")
add_footer(slide, pg, TOTAL_PAGES)

tf = add_content_box(slide, Inches(0.5), Inches(1.6), Inches(6), Inches(5))
add_bullet(tf, "\u76ee\u6807", font_size=20, bold=True, color=DARK_BLUE)
add_bullet(tf,
           "\u5b9e\u73b0 Guest MMIO \u8bfb\u5199 \u2192 TLP \u8f6c\u53d1 \u2192 VCS EP \u5bc4\u5b58\u5668",
           level=1, font_size=16)
add_bullet(tf, "", font_size=10)
add_bullet(tf, "\u5b9e\u73b0\u8981\u70b9", font_size=20, bold=True,
           color=DARK_BLUE)
add_bullet(tf,
           "cosim_mmio_write: \u6784\u9020 MWr TLP, \u586b\u5145 addr/data/len",
           level=1, font_size=16)
add_bullet(tf,
           "cosim_mmio_read: \u6784\u9020 MRd TLP, \u7b49\u5f85 CplD \u8fd4\u56de",
           level=1, font_size=16)
add_bullet(tf,
           "bridge_send_tlp: \u5199\u5165 SHM Ring Buffer",
           level=1, font_size=16)
add_bullet(tf,
           "VCS cosim_recv_tlp: DPI \u8bfb\u53d6 TLP, \u5206\u6d3e\u5230 EP Stub",
           level=1, font_size=16)
add_bullet(tf, "", font_size=10)
add_bullet(tf, "\u6d4b\u8bd5\u65b9\u6cd5", font_size=20, bold=True,
           color=DARK_BLUE)
add_bullet(tf,
           "devmem_test.sh: \u4f7f\u7528 devmem2 \u8bfb\u5199 BAR0 \u5bc4\u5b58\u5668",
           level=1, font_size=16)
add_bullet(tf,
           "\u9a8c\u8bc1: \u5199\u5165\u503c == \u8bfb\u56de\u503c",
           level=1, font_size=16)

add_code_block(slide, Inches(7), Inches(1.8), Inches(5.8), Inches(4.5),
               """# devmem_test.sh core logic
BAR0=$(get_bar0_addr)

# Test 1: write 0xDEADBEEF -> read back
devmem2 $BAR0 w 0xDEADBEEF
val=$(devmem2 $BAR0 w)
assert_eq $val 0xDEADBEEF

# Test 2: write 0x12345678 -> read back
devmem2 $((BAR0+4)) w 0x12345678
val=$(devmem2 $((BAR0+4)) w)
assert_eq $val 0x12345678

# Test 3/4: different offsets
...""")

# ======================================================================
# SLIDE 8 -- Phase 0 Test Results
# ======================================================================
slide, pg = new_slide()
add_title_bar(slide,
              "Phase 0: \u6d4b\u8bd5\u7ed3\u679c",
              "devmem_test 4/4 PASS")
add_footer(slide, pg, TOTAL_PAGES)

data = [
    ["\u6d4b\u8bd5\u9879", "\u64cd\u4f5c", "\u671f\u671b\u503c",
     "\u5b9e\u9645\u503c", "\u7ed3\u679c"],
    ["Test 1", "W/R 0xDEADBEEF @ +0x00", "0xDEADBEEF", "0xDEADBEEF", "PASS"],
    ["Test 2", "W/R 0x12345678 @ +0x04", "0x12345678", "0x12345678", "PASS"],
    ["Test 3", "W/R 0xAAAA5555 @ +0x08", "0xAAAA5555", "0xAAAA5555", "PASS"],
    ["Test 4", "W/R 0x00000000 @ +0x0C", "0x00000000", "0x00000000", "PASS"],
]
add_table(slide, Inches(1), Inches(2.0), Inches(11), Inches(2.5), 5, 5, data,
          col_widths=[Inches(1.5), Inches(3.5), Inches(2), Inches(2),
                      Inches(2)])

shape = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.5), Inches(5.0), Inches(6),
    Inches(1.2)
)
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)
shape.line.color.rgb = GREEN
shape.line.width = Pt(3)
tf = shape.text_frame
p = tf.paragraphs[0]
p.text = "Phase 0 \u7ed3\u679c: 4/4 ALL PASS"
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = GREEN
p.font.name = CN_FONT
p.alignment = PP_ALIGN.CENTER

# ======================================================================
# SLIDE 9 -- Phase 1: Config Space Forwarding
# ======================================================================
slide, pg = new_slide()
add_title_bar(slide,
              "Phase 1: Config Space \u8f6c\u53d1",
              "CfgRd / CfgWr TLP")
add_footer(slide, pg, TOTAL_PAGES)

tf = add_content_box(slide, Inches(0.5), Inches(1.6), Inches(6), Inches(5))
add_bullet(tf, "\u76ee\u6807", font_size=20, bold=True, color=DARK_BLUE)
add_bullet(tf,
           "\u62e6\u622a Guest PCI Config \u8bfb\u5199, \u8f6c\u53d1\u4e3a CfgRd/CfgWr TLP",
           level=1, font_size=16)
add_bullet(tf, "", font_size=10)
add_bullet(tf, "QEMU \u4fa7\u5b9e\u73b0", font_size=20, bold=True,
           color=DARK_BLUE)
add_bullet(tf,
           "cosim_config_read: \u6ce8\u518c\u4e3a pci_device config_read \u56de\u8c03",
           level=1, font_size=16)
add_bullet(tf,
           "cosim_config_write: \u6ce8\u518c\u4e3a config_write \u56de\u8c03",
           level=1, font_size=16)
add_bullet(tf,
           "\u6784\u9020 Type 0 CfgRd/CfgWr TLP (Bus/Dev/Fn + Register)",
           level=1, font_size=16)
add_bullet(tf, "", font_size=10)
add_bullet(tf, "VCS \u4fa7\u5b9e\u73b0", font_size=20, bold=True,
           color=DARK_BLUE)
add_bullet(tf,
           "cfg_space[0:255]: 256 \u5b57\u8282 PCI \u914d\u7f6e\u7a7a\u95f4\u6570\u7ec4",
           level=1, font_size=16)
add_bullet(tf,
           "CfgRd: \u4ece cfg_space \u8bfb\u53d6 dword, \u8fd4\u56de CplD",
           level=1, font_size=16)
add_bullet(tf,
           "CfgWr: \u5199\u5165 cfg_space, \u53d1\u9001 Cpl",
           level=1, font_size=16)

add_code_block(slide, Inches(7), Inches(1.8), Inches(5.8), Inches(4.5),
               """PCI Config Space (Type 0):
+----------------------------------+
| 00h: Vendor ID  |  Device ID     |
| 04h: Command    |  Status        |
| 08h: Rev | Class Code            |
| 0Ch: Cache| Lat | HdrTyp| BIST  |
| 10h: BAR0                        |
| 14h: BAR1 (disabled)             |
| ...                              |
| 34h: Cap Pointer -> 0x40         |
| 3Ch: IntLine|IntPin|MinGnt|MaxL  |
| 40h: Virtio Common Cap           |
| 50h: Virtio Notify Cap           |
| 60h: Virtio ISR Cap              |
| 70h: Virtio Device Cap           |
+----------------------------------+""")

# ======================================================================
# SLIDE 10 -- Phase 1 Issues & Fixes
# ======================================================================
slide, pg = new_slide()
add_title_bar(slide,
              "Phase 1: \u9047\u5230\u7684\u95ee\u9898\u4e0e\u89e3\u51b3")
add_footer(slide, pg, TOTAL_PAGES)

tf = add_content_box(slide, Inches(0.5), Inches(1.6), Inches(12), Inches(1.4))
add_bullet(tf,
           "\u95ee\u9898 1: Header Type \u8bfb\u53d6\u8fd4\u56de 0x10 (\u5e94\u4e3a 0x00)",
           font_size=18, bold=True, color=RED)
add_bullet(tf,
           "\u539f\u56e0: CfgRd \u8fd4\u56de\u6574\u4e2a dword, \u4f46\u672a\u8003\u8651 byte offset\u3002offset=0x0E \u5728 dword 0x0C \u5185, \u9700\u8981\u53f3\u79fb 2 \u5b57\u8282",
           level=1, font_size=15)
add_bullet(tf,
           "\u4fee\u590d: \u5bf9\u9f50\u5230 dword \u8fb9\u754c, \u8bfb\u53d6\u540e\u6309 (offset & 3) \u505a\u5b57\u8282\u79fb\u4f4d: data >>= (byte_offset * 8)",
           level=1, font_size=15, color=GREEN)

tf2 = add_content_box(slide, Inches(0.5), Inches(3.3), Inches(12), Inches(1.4))
add_bullet(tf2,
           "\u95ee\u9898 2: CfgWr \u5199\u5165\u6570\u636e\u5168\u4e3a\u96f6",
           font_size=18, bold=True, color=RED)
add_bullet(tf2,
           "\u539f\u56e0: VCS \u7aef\u7528 tlp_len[1:0] \u63d0\u53d6\u957f\u5ea6, \u4f46 len=4 \u65f6 [1:0] = 0 \u2192 \u5199\u5165 0 \u5b57\u8282",
           level=1, font_size=15)
add_bullet(tf2,
           "\u4fee\u590d: \u6539\u7528 tlp_len[2:0] (3 bit), \u6b63\u786e\u8868\u793a 1-4 \u5b57\u8282\u957f\u5ea6",
           level=1, font_size=15, color=GREEN)

tf3 = add_content_box(slide, Inches(0.5), Inches(5.0), Inches(12), Inches(1.4))
add_bullet(tf3,
           "\u95ee\u9898 3: BAR0 \u5730\u5740\u59cb\u7ec8\u4e3a 0x00000000",
           font_size=18, bold=True, color=RED)
add_bullet(tf3,
           "\u539f\u56e0: Guest \u5199 0xFFFFFFFF \u5230 BAR0 \u505a sizing, \u4f46 VCS \u7aef\u76f4\u63a5\u5b58\u50a8 \u2192 \u8bfb\u56de 0xFFFFFFFF (\u65e0\u6548)",
           level=1, font_size=15)
add_bullet(tf3,
           "\u4fee\u590d: \u6dfb\u52a0 BAR0_SIZE_MASK \u903b\u8f91, \u5199\u5165\u65f6\u4fdd\u5b58\u539f\u503c, sizing \u8bfb\u53d6\u8fd4\u56de ~(SIZE-1) | type bits",
           level=1, font_size=15, color=GREEN)

# ======================================================================
# SLIDE 11 -- Phase 1 Test Results
# ======================================================================
slide, pg = new_slide()
add_title_bar(slide,
              "Phase 1: \u6d4b\u8bd5\u7ed3\u679c",
              "Config Space \u8bfb\u5199\u9a8c\u8bc1")
add_footer(slide, pg, TOTAL_PAGES)

data = [
    ["\u6d4b\u8bd5\u9879", "\u64cd\u4f5c", "\u671f\u671b\u503c",
     "\u7ed3\u679c"],
    ["Vendor ID", "CfgRd 0x00", "0x1AF4 (Virtio)", "PASS"],
    ["Device ID", "CfgRd 0x02", "0x1041 (net)", "PASS"],
    ["Header Type", "CfgRd 0x0E", "0x00 (Type 0)", "PASS"],
    ["BAR0 Sizing", "Write 0xFFF.. -> Read", "~(SIZE-1)", "PASS"],
    ["BAR0 Assign", "Write base addr", "Correct addr", "PASS"],
    ["Cap Pointer", "CfgRd 0x34", "0x40", "PASS"],
    ["Command Reg", "CfgWr 0x04", "Memory Enable", "PASS"],
]
add_table(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(3.5), 8, 4,
          data, col_widths=[Inches(2.5), Inches(3), Inches(2.5), Inches(2)])

shape = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.5), Inches(5.7), Inches(6),
    Inches(1.0)
)
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)
shape.line.color.rgb = GREEN
shape.line.width = Pt(3)
tf = shape.text_frame
p = tf.paragraphs[0]
p.text = "Phase 1 \u7ed3\u679c: Config Space ALL PASS"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = GREEN
p.font.name = CN_FONT
p.alignment = PP_ALIGN.CENTER

# ======================================================================
# SLIDE 12 -- Phase 2: Virtio Register Set
# ======================================================================
slide, pg = new_slide()
add_title_bar(slide,
              "Phase 2: Virtio \u5bc4\u5b58\u5668\u96c6",
              "VCS EP Stub \u6dfb\u52a0 Virtio-PCI \u5bc4\u5b58\u5668")
add_footer(slide, pg, TOTAL_PAGES)

tf = add_content_box(slide, Inches(0.5), Inches(1.6), Inches(5.5), Inches(5))
add_bullet(tf, "Virtio PCI Capabilities (Config Space)", font_size=18,
           bold=True, color=DARK_BLUE)
add_bullet(tf, "Cap @ 0x40: Common Cfg (BAR0 + 0x1000)", level=1,
           font_size=15)
add_bullet(tf, "Cap @ 0x50: Notify (BAR0 + 0x2000)", level=1, font_size=15)
add_bullet(tf, "Cap @ 0x60: ISR (BAR0 + 0x3000)", level=1, font_size=15)
add_bullet(tf, "Cap @ 0x70: Device Cfg (BAR0 + 0x4000)", level=1,
           font_size=15)
add_bullet(tf, "", font_size=10)
add_bullet(tf, "Common Cfg \u5bc4\u5b58\u5668 (BAR0+0x1000)", font_size=18,
           bold=True, color=DARK_BLUE)
add_bullet(tf, "+0x00: device_feature_select", level=1, font_size=15)
add_bullet(tf, "+0x04: device_feature", level=1, font_size=15)
add_bullet(tf, "+0x08: driver_feature_select", level=1, font_size=15)
add_bullet(tf, "+0x0C: driver_feature", level=1, font_size=15)
add_bullet(tf, "+0x14: device_status", level=1, font_size=15)
add_bullet(tf, "+0x16: num_queues", level=1, font_size=15)
add_bullet(tf, "+0x1E: queue_size", level=1, font_size=15)
add_bullet(tf, "+0x20: queue_msix_vector", level=1, font_size=15)
add_bullet(tf, "+0x22: queue_enable", level=1, font_size=15)

tf2 = add_content_box(slide, Inches(6.5), Inches(1.6), Inches(6), Inches(5))
add_bullet(tf2, "Device Cfg \u5bc4\u5b58\u5668 (BAR0+0x4000)", font_size=18,
           bold=True, color=DARK_BLUE)
add_bullet(tf2, "+0x00~0x05: MAC Address (6 bytes)", level=1, font_size=15)
add_bullet(tf2, "+0x06: status (uint16)", level=1, font_size=15)
add_bullet(tf2, "+0x08: max_virtqueue_pairs (uint16)", level=1, font_size=15)
add_bullet(tf2, "", font_size=10)
add_bullet(tf2, "Notify \u5bc4\u5b58\u5668 (BAR0+0x2000)", font_size=18,
           bold=True, color=DARK_BLUE)
add_bullet(tf2, "\u6bcf\u4e2a queue \u4e00\u4e2a notify offset", level=1,
           font_size=15)
add_bullet(tf2, "notify_off_multiplier = 2", level=1, font_size=15)
add_bullet(tf2, "", font_size=10)
add_bullet(tf2, "ISR \u5bc4\u5b58\u5668 (BAR0+0x3000)", font_size=18,
           bold=True, color=DARK_BLUE)
add_bullet(tf2, "+0x00: ISR status (read-clear)", level=1, font_size=15)

# ======================================================================
# SLIDE 13 -- Phase 2 Issues & Fixes
# ======================================================================
slide, pg = new_slide()
add_title_bar(slide,
              "Phase 2: \u9047\u5230\u7684\u95ee\u9898\u4e0e\u89e3\u51b3")
add_footer(slide, pg, TOTAL_PAGES)

tf = add_content_box(slide, Inches(0.5), Inches(1.6), Inches(12), Inches(2.5))
add_bullet(tf, "\u95ee\u9898: MAC \u5730\u5740\u5b57\u8282\u5e8f\u53cd\u8f6c",
           font_size=20, bold=True, color=RED)
add_bullet(tf,
           "\u73b0\u8c61: Guest \u8bfb\u53d6 MAC \u5f97\u5230 01:00:EF:BE:AD:DE (\u5e94\u4e3a DE:AD:BE:EF:00:01)",
           level=0, font_size=16)
add_bullet(tf,
           "\u6839\u56e0: MMIO \u8bfb\u53d6\u8fd4\u56de 32-bit dword, \u4f46 MAC \u4e3a 6 \u5b57\u8282\u6309\u5b57\u8282\u5bfb\u5740",
           level=0, font_size=16)
add_bullet(tf,
           "VCS \u7aef\u5c06 MAC \u6309 big-endian \u6392\u5217\u5728 dword \u4e2d, \u4f46 MMIO \u901a\u8def\u4e3a little-endian",
           level=0, font_size=16)

tf2 = add_content_box(slide, Inches(0.5), Inches(4.0), Inches(12), Inches(1.5))
add_bullet(tf2, "\u4fee\u590d\u65b9\u6848:", font_size=20, bold=True,
           color=GREEN)
add_bullet(tf2,
           "vio_devcfg_read \u51fd\u6570\u4e2d\u6309 LE \u89c4\u5219\u91cd\u6392\u5b57\u8282:",
           level=0, font_size=16)
add_bullet(tf2,
           "MAC[0] -> byte0 (bits[7:0]), MAC[1] -> byte1 (bits[15:8]), ...",
           level=1, font_size=15)

add_code_block(slide, Inches(1), Inches(5.5), Inches(11), Inches(1.2),
               """// vio_devcfg_read: MAC LE reorder
data = {mac[offset+3], mac[offset+2], mac[offset+1], mac[offset]};  // LE byte order""")

# ======================================================================
# SLIDE 14 -- Phase 2 Test Results
# ======================================================================
slide, pg = new_slide()
add_title_bar(slide,
              "Phase 2: \u6d4b\u8bd5\u7ed3\u679c",
              "virtio_reg_test 17/17 ALL PASS")
add_footer(slide, pg, TOTAL_PAGES)

data = [
    ["#", "\u6d4b\u8bd5\u9879", "\u671f\u671b", "\u7ed3\u679c"],
    ["1", "device_feature_select W/R", "0/1", "PASS"],
    ["2", "device_feature (page0)", "Features[31:0]", "PASS"],
    ["3", "device_feature (page1)", "Features[63:32]", "PASS"],
    ["4", "driver_feature_select W/R", "0/1", "PASS"],
    ["5", "driver_feature W/R", "Negotiated", "PASS"],
    ["6", "device_status W/R", "0x01->0x03->0x0F", "PASS"],
    ["7", "num_queues", "2 (TX+RX)", "PASS"],
    ["8", "queue_select W/R", "0/1", "PASS"],
    ["9", "queue_size", "256", "PASS"],
    ["10-13", "queue desc/avail/used addr", "64-bit W/R", "PASS"],
    ["14", "queue_enable", "0->1", "PASS"],
    ["15", "MAC addr (6 bytes)", "DE:AD:BE:EF:00:01", "PASS"],
    ["16", "ISR read-clear", "0x01->0x00", "PASS"],
    ["17", "notify queue kick", "Queue index", "PASS"],
]
add_table(slide, Inches(1), Inches(1.8), Inches(11), Inches(5.0), 15, 4,
          data, col_widths=[Inches(1), Inches(4), Inches(3), Inches(3)])

# ======================================================================
# SLIDE 15 -- Phase 2b: Virtio Driver Binding
# ======================================================================
slide, pg = new_slide()
add_title_bar(slide,
              "Phase 2b: Virtio \u9a71\u52a8\u7ed1\u5b9a",
              "virtio_net.ko \u5185\u6838\u6a21\u5757\u52a0\u8f7d")
add_footer(slide, pg, TOTAL_PAGES)

tf = add_content_box(slide, Inches(0.5), Inches(1.6), Inches(6), Inches(5))
add_bullet(tf, "\u76ee\u6807", font_size=20, bold=True, color=DARK_BLUE)
add_bullet(tf,
           "\u8ba9 Linux \u5185\u6838 virtio_net \u9a71\u52a8\u6210\u529f probe \u5e76\u521b\u5efa eth0",
           level=1, font_size=16)
add_bullet(tf, "", font_size=10)
add_bullet(tf, "\u5173\u952e\u914d\u7f6e", font_size=20, bold=True,
           color=DARK_BLUE)
add_bullet(tf, "Vendor ID: 0x1AF4 (Red Hat / Virtio)", level=1, font_size=16)
add_bullet(tf, "Device ID: 0x1041 (virtio-net, transitional)", level=1,
           font_size=16)
add_bullet(tf, "Subsystem Vendor: 0x1AF4", level=1, font_size=16)
add_bullet(tf, "Subsystem Device: 0x0001 (network)", level=1, font_size=16)
add_bullet(tf, "Revision: 0x01 (virtio 1.0+)", level=1, font_size=16)
add_bullet(tf, "", font_size=10)
add_bullet(tf, "\u9a71\u52a8\u52a0\u8f7d\u6d41\u7a0b", font_size=20,
           bold=True, color=DARK_BLUE)
add_bullet(tf, "1. modprobe virtio_pci", level=1, font_size=16)
add_bullet(tf, "2. PCI \u679a\u4e3e\u5339\u914d Vendor/Device ID", level=1,
           font_size=16)
add_bullet(tf, "3. virtio_pci_probe -> \u8bfb\u53d6 capabilities", level=1,
           font_size=16)
add_bullet(tf, "4. virtio_net probe -> \u534f\u5546 features", level=1,
           font_size=16)
add_bullet(tf, "5. \u521b\u5efa eth0 \u7f51\u5361\u63a5\u53e3", level=1,
           font_size=16)

flow_steps = [
    "PCI \u679a\u4e3e\u53d1\u73b0\u8bbe\u5907",
    "virtio_pci_probe",
    "\u8bfb\u53d6 Capabilities",
    "Feature \u534f\u5546",
    "Queue \u914d\u7f6e",
    "virtio_net probe",
    "eth0 \u521b\u5efa\u6210\u529f"
]
y_start = Inches(1.8)
for i, step in enumerate(flow_steps):
    add_diagram_box(slide, Inches(8), y_start + Inches(i * 0.7),
                    Inches(4.5), Inches(0.5), step, font_size=14)
    if i < len(flow_steps) - 1:
        add_down_arrow(slide, Inches(10.1),
                       y_start + Inches(i * 0.7 + 0.5),
                       Inches(0.3), Inches(0.18))

# ======================================================================
# SLIDE 16 -- Phase 2b Issues & Fixes
# ======================================================================
slide, pg = new_slide()
add_title_bar(slide,
              "Phase 2b: \u9047\u5230\u7684\u95ee\u9898\u4e0e\u89e3\u51b3")
add_footer(slide, pg, TOTAL_PAGES)

tf = add_content_box(slide, Inches(0.5), Inches(1.5), Inches(12), Inches(1.5))
add_bullet(tf,
           "\u95ee\u9898 1: virtio_net.ko \u6a21\u5757\u7f3a\u5931",
           font_size=18, bold=True, color=RED)
add_bullet(tf,
           "\u73b0\u8c61: modprobe virtio_net \u62a5\u9519 module not found",
           level=0, font_size=15)
add_bullet(tf,
           "\u6839\u56e0: QEMU \u4f7f\u7528\u7684 rootfs \u5185\u6838 (6.6.x) \u672a\u5305\u542b virtio_net \u6a21\u5757",
           level=0, font_size=15)
add_bullet(tf,
           "\u4fee\u590d: \u4e0b\u8f7d\u5339\u914d\u7248\u672c linux-modules-6.6.110-virt, \u89e3\u538b\u5230 rootfs \u7684 /lib/modules/",
           level=0, font_size=15, color=GREEN)

tf2 = add_content_box(slide, Inches(0.5), Inches(3.3), Inches(12),
                      Inches(1.7))
add_bullet(tf2,
           "\u95ee\u9898 2: probe failed -22 (EINVAL)",
           font_size=18, bold=True, color=RED)
add_bullet(tf2,
           "\u73b0\u8c61: virtio_pci_probe \u8fd4\u56de -EINVAL, dmesg \u663e\u793a IRQ \u5206\u914d\u5931\u8d25",
           level=0, font_size=15)
add_bullet(tf2,
           "\u6839\u56e0: Config Space INT_PIN = 0 -> \u5185\u6838\u8ba4\u4e3a\u8bbe\u5907\u4e0d\u652f\u6301\u4e2d\u65ad -> IRQ=0",
           level=0, font_size=15)
add_bullet(tf2,
           "\u4fee\u590d: \u8bbe\u7f6e cfg_space[0x3D] = 0x01 (INT_PIN = INTA), \u5185\u6838\u6b63\u786e\u5206\u914d IRQ",
           level=0, font_size=15, color=GREEN)

tf3 = add_content_box(slide, Inches(0.5), Inches(5.3), Inches(12),
                      Inches(1.5))
add_bullet(tf3,
           "\u95ee\u9898 3: BAR1-5 \u88ab\u8bef\u5206\u914d\u5730\u5740\u7a7a\u95f4",
           font_size=18, bold=True, color=RED)
add_bullet(tf3,
           "\u73b0\u8c61: \u5185\u6838\u4e3a BAR1-5 \u5206\u914d\u4e86\u5730\u5740, \u4f46\u5b9e\u9645\u4e0d\u5b58\u5728 -> \u8bbf\u95ee\u5f02\u5e38",
           level=0, font_size=15)
add_bullet(tf3,
           "\u6839\u56e0: CfgWr \u5bf9 BAR1-5 \u5199 0xFFFFFFFF \u540e, VCS \u7aef\u672a\u8fd4\u56de 0 (\u5e94\u6307\u793a\u4e0d\u53ef\u7528)",
           level=0, font_size=15)
add_bullet(tf3,
           "\u4fee\u590d: CfgRd BAR1-5 \u59cb\u7ec8\u8fd4\u56de 0x00000000, \u8868\u793a\u8fd9\u4e9b BAR \u672a\u5b9e\u73b0",
           level=0, font_size=15, color=GREEN)

# ======================================================================
# SECTION DIVIDER -- Part III
# ======================================================================
slide, pg = new_slide()
add_section_divider(slide, "III",
                    "\u7b2c\u4e09\u90e8\u5206: Phase 2b \u6700\u7ec8\u7ed3\u679c",
                    "\u9a71\u52a8\u7ed1\u5b9a\u6210\u529f + \u5b8c\u6574\u9a8c\u8bc1")
add_footer(slide, pg, TOTAL_PAGES)

# ======================================================================
# SLIDE 18 -- Phase 2b Final Test Results
# ======================================================================
slide, pg = new_slide()
add_title_bar(slide,
              "Phase 2b: \u6700\u7ec8\u6d4b\u8bd5\u7ed3\u679c",
              "virtio_net \u9a71\u52a8\u7ed1\u5b9a\u6210\u529f")
add_footer(slide, pg, TOTAL_PAGES)

results = [
    ("device_status", "0x0F (DRIVER_OK)", GREEN),
    ("\u7f51\u5361\u63a5\u53e3", "eth0 \u5df2\u521b\u5efa", GREEN),
    ("MAC \u5730\u5740", "DE:AD:BE:EF:00:01", GREEN),
    ("\u9a71\u52a8\u6a21\u5757", "virtio_net.ko loaded", GREEN),
    ("PCI \u8bbe\u5907", "1AF4:1041 (virtio-net)", GREEN),
]

for i, (label, value, color) in enumerate(results):
    y = Inches(1.8) + Inches(i * 0.8)
    add_diagram_box(slide, Inches(1), y, Inches(3), Inches(0.6), label,
                    font_size=16)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.5), y, Inches(5), Inches(0.6)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)
    shape.line.color.rgb = color
    shape.line.width = Pt(2)
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = value
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = color
    p.font.name = CODE_FONT
    p.alignment = PP_ALIGN.CENTER

add_code_block(slide, Inches(1), Inches(5.8), Inches(11), Inches(1.2),
               """# dmesg key output
[  2.345] virtio_net virtio0: device_status = 0x0F (DRIVER_OK)
[  2.346] virtio_net virtio0 eth0: MAC DE:AD:BE:EF:00:01, MTU 1500""")

# ======================================================================
# SLIDE 19 -- VCS Full Negotiation Flow
# ======================================================================
slide, pg = new_slide()
add_title_bar(slide,
              "VCS \u4fa7: Virtio \u5b8c\u6574\u534f\u5546\u6d41\u7a0b")
add_footer(slide, pg, TOTAL_PAGES)

steps = [
    ("1. Reset", "device_status = 0x00",
     "\u9a71\u52a8\u91cd\u7f6e\u8bbe\u5907"),
    ("2. ACK", "device_status |= 0x01",
     "\u9a71\u52a8\u786e\u8ba4\u53d1\u73b0\u8bbe\u5907"),
    ("3. DRIVER", "device_status |= 0x02",
     "\u9a71\u52a8\u5df2\u52a0\u8f7d"),
    ("4. Features",
     "\u8bfb\u53d6 device_feature\n\u5199\u5165 driver_feature",
     "\u534f\u5546\u8bbe\u5907\u7279\u6027"),
    ("5. FEATURES_OK", "device_status |= 0x08",
     "\u7279\u6027\u534f\u5546\u5b8c\u6210"),
    ("6. Queue Setup",
     "\u914d\u7f6e queue_size\ndesc/avail/used addr\nqueue_enable = 1",
     "\u521d\u59cb\u5316 Virtqueue"),
    ("7. DRIVER_OK", "device_status |= 0x04",
     "\u9a71\u52a8\u5c31\u7eea, \u8bbe\u5907\u53ef\u7528"),
]

for i, (step, detail, desc) in enumerate(steps):
    y = Inches(1.6) + Inches(i * 0.75)
    add_diagram_box(slide, Inches(0.5), y, Inches(2.5), Inches(0.6), step,
                    bg_color=RGBColor(0xE3, 0xF2, 0xFD), font_size=14)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.3), y, Inches(4.5), Inches(0.6)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_GRAY
    shape.line.color.rgb = MEDIUM_GRAY
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = detail
    p.font.size = Pt(12)
    p.font.name = CODE_FONT
    p.font.color.rgb = DARK_GRAY
    p.alignment = PP_ALIGN.CENTER
    tf2 = add_content_box(slide, Inches(8.2), y, Inches(4.5), Inches(0.6))
    p2 = tf2.paragraphs[0]
    p2.text = desc
    p2.font.size = Pt(14)
    p2.font.name = CN_FONT
    p2.font.color.rgb = DARK_GRAY

    if i < len(steps) - 1:
        add_down_arrow(slide, Inches(1.6), y + Inches(0.6), Inches(0.25),
                       Inches(0.12))

# ======================================================================
# SLIDE 20 -- Config Space Hex Dump
# ======================================================================
slide, pg = new_slide()
add_title_bar(slide, "Config Space Hex Dump",
              "\u4fee\u590d\u540e\u7684\u6b63\u786e\u914d\u7f6e")
add_footer(slide, pg, TOTAL_PAGES)

add_code_block(slide, Inches(0.5), Inches(1.6), Inches(12.3), Inches(5.2),
               """PCI Config Space (256 bytes) - after fixes:

Offset  00 01 02 03  04 05 06 07  08 09 0A 0B  0C 0D 0E 0F
------  ----------- ----------- ----------- -----------
  00h:  F4 1A 41 10  06 01 10 00  01 00 00 02  00 00 00 00   <- Vendor=1AF4 Device=1041
  10h:  xx xx xx x4  00 00 00 00  00 00 00 00  00 00 00 00   <- BAR0=actual addr BAR1-5=0
  20h:  00 00 00 00  00 00 00 00  00 00 00 00  F4 1A 01 00   <- SubVendor=1AF4 SubDev=0001
  30h:  00 00 00 00  40 00 00 00  00 00 00 00  0B 01 00 00   <- CapPtr=0x40 IntPin=INTA(0x01)

  40h:  09 50 01 01  04 00 00 10  00 10 00 00  00 00 00 00   <- Virtio Common Cfg Cap
  50h:  09 60 02 01  04 00 00 20  00 10 00 00  02 00 00 00   <- Virtio Notify Cap (mult=2)
  60h:  09 70 03 01  04 00 00 30  00 10 00 00  00 00 00 00   <- Virtio ISR Cap
  70h:  09 00 04 01  04 00 00 40  00 10 00 00  00 00 00 00   <- Virtio Device Cfg Cap (last)

Key fix annotations:
  [0x0E] Header Type = 0x00 (Type 0 endpoint)          <- Phase 1 fix
  [0x10] BAR0 = actual assigned address (not 0)         <- BAR sizing fix
  [0x14-0x24] BAR1-5 = 0x00000000 (not implemented)    <- Phase 2b fix
  [0x3D] INT_PIN = 0x01 (INTA)                          <- Phase 2b fix""")

# ======================================================================
# SLIDE 21 -- DMA Path Verification
# ======================================================================
slide, pg = new_slide()
add_title_bar(slide, "DMA \u901a\u8def\u9a8c\u8bc1")
add_footer(slide, pg, TOTAL_PAGES)

tf = add_content_box(slide, Inches(0.5), Inches(1.6), Inches(6), Inches(3))
add_bullet(tf, "DMA Read \u6d41\u7a0b", font_size=20, bold=True,
           color=DARK_BLUE)
add_bullet(tf,
           "1. VCS EP \u53d1\u8d77 DMA Read Request (\u5730\u5740 + \u957f\u5ea6)",
           level=1, font_size=15)
add_bullet(tf,
           "2. bridge \u63a5\u6536\u8bf7\u6c42, \u8c03\u7528 QEMU address_space_read",
           level=1, font_size=15)
add_bullet(tf,
           "3. \u4ece Guest \u7269\u7406\u5185\u5b58\u8bfb\u53d6\u6570\u636e",
           level=1, font_size=15)
add_bullet(tf,
           "4. \u901a\u8fc7 Ring Buffer \u8fd4\u56de DMA Read Completion",
           level=1, font_size=15)
add_bullet(tf, "", font_size=10)
add_bullet(tf, "DMA Write \u6d41\u7a0b", font_size=20, bold=True,
           color=DARK_BLUE)
add_bullet(tf,
           "1. VCS EP \u53d1\u8d77 DMA Write (\u5730\u5740 + \u6570\u636e)",
           level=1, font_size=15)
add_bullet(tf,
           "2. bridge \u63a5\u6536, \u8c03\u7528 address_space_write",
           level=1, font_size=15)
add_bullet(tf,
           "3. \u5199\u5165 Guest \u7269\u7406\u5185\u5b58",
           level=1, font_size=15)

data = [
    ["DMA \u6d4b\u8bd5\u9879", "\u65b9\u5411", "\u5927\u5c0f",
     "\u7ed3\u679c"],
    ["DMA Read \u5355 dword", "VCS->Guest", "4B", "PASS"],
    ["DMA Read \u591a dword", "VCS->Guest", "64B", "PASS"],
    ["DMA Write \u5355 dword", "VCS->Guest", "4B", "PASS"],
    ["DMA Write \u591a dword", "VCS->Guest", "64B", "PASS"],
    ["DMA \u5730\u5740\u5bf9\u9f50", "\u53cc\u5411", "Various", "PASS"],
]
add_table(slide, Inches(6.8), Inches(1.8), Inches(6), Inches(2.8), 6, 4,
          data, col_widths=[Inches(2), Inches(1.2), Inches(1), Inches(1.8)])

shape = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(5.2), Inches(5.5),
    Inches(0.8)
)
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)
shape.line.color.rgb = GREEN
shape.line.width = Pt(3)
tf2 = shape.text_frame
p = tf2.paragraphs[0]
p.text = "DMA \u901a\u8def: ALL PASS"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = GREEN
p.font.name = CN_FONT
p.alignment = PP_ALIGN.CENTER

# ======================================================================
# SLIDE 22 -- NIC TX Simulation
# ======================================================================
slide, pg = new_slide()
add_title_bar(slide,
              "NIC TX \u4eff\u771f: \u6570\u636e\u5305\u53d1\u9001\u9a8c\u8bc1")
add_footer(slide, pg, TOTAL_PAGES)

tf = add_content_box(slide, Inches(0.5), Inches(1.6), Inches(6), Inches(4.5))
add_bullet(tf, "TX \u53d1\u9001\u6d41\u7a0b", font_size=20, bold=True,
           color=DARK_BLUE)
add_bullet(tf,
           "1. Guest \u6784\u9020\u7f51\u7edc\u5305, \u5199\u5165 TX Virtqueue descriptor",
           level=1, font_size=15)
add_bullet(tf,
           "2. \u66f4\u65b0 avail ring, \u5199 notify \u5bc4\u5b58\u5668 kick",
           level=1, font_size=15)
add_bullet(tf,
           "3. VCS EP \u68c0\u6d4b notify, DMA \u8bfb\u53d6 descriptor",
           level=1, font_size=15)
add_bullet(tf,
           "4. DMA \u8bfb\u53d6\u6570\u636e\u5305\u5185\u5bb9",
           level=1, font_size=15)
add_bullet(tf,
           "5. \u5199\u5165 used ring, \u89e6\u53d1\u4e2d\u65ad",
           level=1, font_size=15)
add_bullet(tf, "", font_size=10)
add_bullet(tf, "\u9a8c\u8bc1\u7ed3\u679c", font_size=20, bold=True,
           color=DARK_BLUE)
add_bullet(tf,
           "Guest ping -> VCS \u7aef\u6210\u529f\u63a5\u6536 ARP/ICMP \u5305",
           level=1, font_size=15)
add_bullet(tf,
           "\u5305\u5185\u5bb9\u4e0e Guest \u53d1\u9001\u4e00\u81f4",
           level=1, font_size=15)

flow = [
    "Guest: \u6784\u9020 packet",
    "\u5199\u5165 TX desc ring",
    "Notify kick (BAR0+0x2000)",
    "VCS: DMA read desc",
    "VCS: DMA read packet",
    "VCS: Update used ring",
    "\u89e6\u53d1\u4e2d\u65ad -> Guest \u56de\u6536 buffer",
]
for i, step in enumerate(flow):
    add_diagram_box(
        slide, Inches(7.5), Inches(1.6) + Inches(i * 0.7), Inches(5),
        Inches(0.5), step, font_size=13,
        bg_color=(RGBColor(0xE3, 0xF2, 0xFD) if i < 3
                  else RGBColor(0xE8, 0xF5, 0xE9)))
    if i < len(flow) - 1:
        add_down_arrow(slide, Inches(9.85),
                       Inches(2.1) + Inches(i * 0.7),
                       Inches(0.25), Inches(0.15))

# ======================================================================
# SECTION DIVIDER -- Part IV: Phase 3
# ======================================================================
slide, pg = new_slide()
add_section_divider(slide, "IV",
                    "Phase 3: Virtqueue \u5f15\u64ce\u5b9e\u73b0",
                    "VCS \u7aef\u5b8c\u6574\u5904\u7406 Virtqueue \u63cf\u8ff0\u7b26\u94fe + ETH SHM \u8f6c\u53d1")
add_footer(slide, pg, TOTAL_PAGES)

# ======================================================================
# SLIDE 24 -- Phase 3: Architecture
# ======================================================================
slide, pg = new_slide()
add_title_bar(slide,
              "Phase 3: Virtqueue \u67b6\u6784\u8bbe\u8ba1",
              "\u65b0\u589e\u6a21\u5757\u4e0e\u6570\u636e\u6d41")
add_footer(slide, pg, TOTAL_PAGES)

tf = add_content_box(slide, Inches(0.5), Inches(1.6), Inches(6), Inches(5))
add_bullet(tf, "\u65b0\u589e\u6587\u4ef6", font_size=20, bold=True, color=DARK_BLUE)
add_bullet(tf, "virtqueue_dma.h / virtqueue_dma.c  \u2014  VQ \u5904\u7406\u5f15\u64ce",
           level=1, font_size=15)
add_bullet(tf, "", font_size=10)
add_bullet(tf, "\u4fee\u6539\u6587\u4ef6", font_size=20, bold=True, color=DARK_BLUE)
add_bullet(tf, "bridge_vcs.c  \u2014  \u65b0\u589e DMA \u4efb\u610f\u957f\u5ea6\u8bfb\u5199 (uint8_t*)",
           level=1, font_size=15)
add_bullet(tf, "eth_mac_dpi.c  \u2014  \u65b0\u589e raw C \u5c42 send/recv \u51fd\u6570",
           level=1, font_size=15)
add_bullet(tf, "pcie_ep_stub.sv  \u2014  \u65b0\u589e notify_valid/notify_queue \u8f93\u51fa",
           level=1, font_size=15)
add_bullet(tf, "tb_top.sv  \u2014  DRIVER_OK \u68c0\u6d4b + VQ \u914d\u7f6e + notify \u5904\u7406",
           level=1, font_size=15)
add_bullet(tf, "", font_size=10)
add_bullet(tf, "\u6838\u5fc3 API", font_size=20, bold=True, color=DARK_BLUE)
add_bullet(tf, "bridge_dma_read_bytes(addr, buf, len)  \u2014  \u4efb\u610f\u957f\u5ea6 DMA \u8bfb",
           level=1, font_size=15)
add_bullet(tf, "bridge_dma_write_bytes(addr, buf, len)  \u2014  \u4efb\u610f\u957f\u5ea6 DMA \u5199",
           level=1, font_size=15)
add_bullet(tf, "vcs_vq_process_tx()  \u2014  TX \u961f\u5217\u5904\u7406",
           level=1, font_size=15)
add_bullet(tf, "vcs_vq_process_rx()  \u2014  RX \u961f\u5217\u5904\u7406",
           level=1, font_size=15)

# Right side: Virtqueue ring structure diagram
add_diagram_box(slide, Inches(7.5), Inches(1.8), Inches(5), Inches(1.2),
                "Descriptor Table (16B/\u6761\u76ee)\naddr[8] + len[4] + flags[2] + next[2]",
                font_size=14)
add_down_arrow(slide, Inches(9.85), Inches(3.05), Inches(0.3), Inches(0.3))
add_diagram_box(slide, Inches(7.5), Inches(3.4), Inches(5), Inches(0.8),
                "Avail Ring\nflags[2] + idx[2] + ring[2*N]", font_size=14,
                bg_color=RGBColor(0xE3, 0xF2, 0xFD))
add_down_arrow(slide, Inches(9.85), Inches(4.25), Inches(0.3), Inches(0.3))
add_diagram_box(slide, Inches(7.5), Inches(4.6), Inches(5), Inches(0.8),
                "Used Ring\nflags[2] + idx[2] + ring[8*N](id+len)",
                font_size=14, bg_color=RGBColor(0xE8, 0xF5, 0xE9))

add_code_block(slide, Inches(7.5), Inches(5.6), Inches(5), Inches(1.3),
               """struct virtq_desc {
    uint64_t addr;   // GPA
    uint32_t len;
    uint16_t flags;  // NEXT=1, WRITE=2
    uint16_t next;
};""")

# ======================================================================
# SLIDE 25 -- Phase 3: TX Flow
# ======================================================================
slide, pg = new_slide()
add_title_bar(slide,
              "Phase 3: TX \u53d1\u9001\u6d41\u7a0b",
              "Guest \u53d1\u5305 \u2192 VCS \u63d0\u53d6 \u2192 ETH SHM")
add_footer(slide, pg, TOTAL_PAGES)

flow = [
    ("Guest: \u6784\u9020\u7f51\u7edc\u5305\n\u5199\u5165 TX descriptor",
     RGBColor(0xE3, 0xF2, 0xFD)),
    ("Guest: \u66f4\u65b0 avail ring idx\n\u5199 notify \u5bc4\u5b58\u5668 (BAR0+0x2000)",
     RGBColor(0xE3, 0xF2, 0xFD)),
    ("EP Stub: \u68c0\u6d4b notify MWr\n\u8f93\u51fa notify_valid + queue_idx",
     RGBColor(0xFF, 0xF8, 0xE1)),
    ("tb_top: \u68c0\u6d4b notify_valid\n\u8c03\u7528 vcs_vq_process_tx()",
     RGBColor(0xFF, 0xF8, 0xE1)),
    ("VQ \u5f15\u64ce: DMA \u8bfb avail ring\n\u83b7\u53d6 head descriptor index",
     RGBColor(0xE8, 0xF5, 0xE9)),
    ("VQ \u5f15\u64ce: DMA \u8bfb descriptor chain\n\u8df3\u8fc7 12B virtio-net header",
     RGBColor(0xE8, 0xF5, 0xE9)),
    ("VQ \u5f15\u64ce: DMA \u8bfb\u53d6\u6570\u636e\u5305\n\u8f6c\u53d1\u5230 ETH SHM",
     RGBColor(0xE8, 0xF5, 0xE9)),
    ("VQ \u5f15\u64ce: \u66f4\u65b0 used ring\nused_idx++ \u901a\u77e5 Guest",
     RGBColor(0xE8, 0xF5, 0xE9)),
]
for i, (step, bg) in enumerate(flow):
    y = Inches(1.5) + Inches(i * 0.7)
    add_diagram_box(slide, Inches(0.5), y, Inches(6.5), Inches(0.55), step,
                    bg_color=bg, font_size=12)
    if i < len(flow) - 1:
        add_down_arrow(slide, Inches(3.6), y + Inches(0.55),
                       Inches(0.25), Inches(0.12))

# Right side: virtio-net header stripping
tf = add_content_box(slide, Inches(7.5), Inches(1.5), Inches(5.3), Inches(2.5))
add_bullet(tf, "Virtio-net Header \u5904\u7406", font_size=18, bold=True,
           color=DARK_BLUE)
add_bullet(tf, "VERSION_1 header = 12 \u5b57\u8282", level=1, font_size=15)
add_bullet(tf, "\u5fc5\u987b\u8df3\u8fc7\u540e\u624d\u662f\u7f51\u7edc\u5305\u6570\u636e", level=1, font_size=15)
add_bullet(tf, "Linux \u53ef\u80fd\u5408\u5e76 header+data \u5728\u5355\u4e2a\u63cf\u8ff0\u7b26", level=1, font_size=15)
add_bullet(tf, "\u4f7f\u7528 hdr_remaining \u8ba1\u6570\u5668\u8ddf\u8e2a\u5269\u4f59\u5934\u90e8\u5b57\u8282",
           level=1, font_size=15)

add_code_block(slide, Inches(7.5), Inches(4.2), Inches(5.3), Inches(2.8),
               """// hdr_remaining = 12 (virtio-net hdr)
if (hdr_remaining > 0) {
    if (buf_len <= hdr_remaining) {
        // \u6574\u4e2a desc \u90fd\u662f header
        hdr_remaining -= buf_len;
        buf_len = 0;
    } else {
        // \u90e8\u5206 header, \u8df3\u8fc7\u540e\u8bfb data
        buf_addr += hdr_remaining;
        buf_len  -= hdr_remaining;
        hdr_remaining = 0;
    }
}""")

# ======================================================================
# SLIDE 26 -- Phase 3: Problems & Fixes
# ======================================================================
slide, pg = new_slide()
add_title_bar(slide,
              "Phase 3: \u9047\u5230\u7684\u95ee\u9898\u4e0e\u89e3\u51b3")
add_footer(slide, pg, TOTAL_PAGES)

tf = add_content_box(slide, Inches(0.5), Inches(1.5), Inches(12), Inches(1.3))
add_bullet(tf,
           "\u95ee\u9898 1: virtio_net.ko \u6a21\u5757\u8def\u5f84\u9519\u8bef",
           font_size=18, bold=True, color=RED)
add_bullet(tf,
           "\u73b0\u8c61: insmod \u62a5\u9519 module not found, device_status \u505c\u5728 0x01",
           level=0, font_size=15)
add_bullet(tf,
           "\u4fee\u590d: \u4f7f\u7528 KVER=$(uname -r) \u52a8\u6001\u62fc\u63a5\u8def\u5f84 /lib/modules/$KVER/kernel/drivers/net/virtio_net.ko",
           level=0, font_size=15, color=GREEN)

tf2 = add_content_box(slide, Inches(0.5), Inches(3.0), Inches(12), Inches(1.3))
add_bullet(tf2,
           "\u95ee\u9898 2: NBA \u65f6\u5e8f\u95ee\u9898 \u2014 notify_valid \u672a\u88ab\u68c0\u6d4b",
           font_size=18, bold=True, color=RED)
add_bullet(tf2,
           "\u73b0\u8c61: EP Stub \u8f93\u51fa notify_valid (NBA), \u4f46 tb_top initial \u5757\u540c\u5468\u671f\u68c0\u67e5\u65f6\u8fd8\u672a\u66f4\u65b0",
           level=0, font_size=15)
add_bullet(tf2,
           "\u4fee\u590d: \u5728\u68c0\u67e5 notify_valid \u524d\u6dfb\u52a0 @(posedge clk) \u7b49\u5f85 NBA \u4f20\u64ad",
           level=0, font_size=15, color=GREEN)

tf3 = add_content_box(slide, Inches(0.5), Inches(4.5), Inches(12), Inches(1.3))
add_bullet(tf3,
           "\u95ee\u9898 3: Virtio-net Header \u672a\u5265\u79bb",
           font_size=18, bold=True, color=RED)
add_bullet(tf3,
           "\u73b0\u8c61: \u8f6c\u53d1 54 \u5b57\u8282/\u5305 (\u5e94\u4e3a 42), \u524d 12 \u5b57\u8282\u4e3a virtio header \u800c\u975e\u7f51\u7edc\u6570\u636e",
           level=0, font_size=15)
add_bullet(tf3,
           "\u4fee\u590d: \u4f7f\u7528 hdr_remaining=12 \u8ba1\u6570\u5668, \u652f\u6301 header \u4e0e data \u5408\u5e76\u5728\u5355\u4e2a\u63cf\u8ff0\u7b26\u7684\u60c5\u51b5",
           level=0, font_size=15, color=GREEN)

tf4 = add_content_box(slide, Inches(0.5), Inches(6.0), Inches(12), Inches(1.0))
add_bullet(tf4,
           "\u95ee\u9898 4: DMA \u4ec5\u652f\u6301 64 \u5b57\u8282 (int[16])",
           font_size=18, bold=True, color=RED)
add_bullet(tf4,
           "\u4fee\u590d: \u65b0\u589e bridge_dma_read/write_bytes() \u2014 \u652f\u6301\u4efb\u610f\u957f\u5ea6 uint8_t* \u7f13\u51b2\u533a",
           level=0, font_size=15, color=GREEN)

# ======================================================================
# SLIDE 27 -- Phase 3: Test Results
# ======================================================================
slide, pg = new_slide()
add_title_bar(slide,
              "Phase 3: \u6d4b\u8bd5\u7ed3\u679c",
              "Virtqueue TX \u8def\u5f84\u5b8c\u6574\u9a8c\u8bc1")
add_footer(slide, pg, TOTAL_PAGES)

data = [
    ["\u6d4b\u8bd5\u9879", "\u8be6\u60c5", "\u7ed3\u679c"],
    ["DRIVER_OK \u68c0\u6d4b", "device_status & 0x04 \u89e6\u53d1 VQ \u914d\u7f6e", "PASS"],
    ["RX Queue \u914d\u7f6e", "desc=0x22c0000 avail=0x22c1000 size=256", "PASS"],
    ["TX Queue \u914d\u7f6e", "desc=0x1fb0000 avail=0x1fb1000 size=256", "PASS"],
    ["TX Pkt #1 (ARP)", "54B desc \u2192 \u5265\u79bb 12B hdr \u2192 42B \u8f6c\u53d1", "PASS"],
    ["TX Pkt #2 (ARP)", "54B desc \u2192 42B \u8f6c\u53d1, used_idx=2", "PASS"],
    ["TX Pkt #3 (ARP)", "54B desc \u2192 42B \u8f6c\u53d1, used_idx=3", "PASS"],
    ["Used Ring \u66f4\u65b0", "used_idx 0\u21921\u21922\u21923 \u6b63\u786e\u9012\u589e", "PASS"],
    ["ETH SHM \u8f6c\u53d1", "\u6240\u6709\u5305\u6210\u529f\u5199\u5165\u5171\u4eab\u5185\u5b58", "PASS"],
]
add_table(slide, Inches(0.5), Inches(1.8), Inches(12.3), Inches(3.8), 9, 3,
          data, col_widths=[Inches(3), Inches(6.3), Inches(3)])

shape = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.5), Inches(5.8), Inches(6),
    Inches(1.0)
)
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)
shape.line.color.rgb = GREEN
shape.line.width = Pt(3)
tf = shape.text_frame
p = tf.paragraphs[0]
p.text = "Phase 3 \u7ed3\u679c: TX \u8def\u5f84 ALL PASS (3/3 \u5305)"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = GREEN
p.font.name = CN_FONT
p.alignment = PP_ALIGN.CENTER

# ======================================================================
# SLIDE 28 -- Phase 3: VCS Log Output
# ======================================================================
slide, pg = new_slide()
add_title_bar(slide,
              "Phase 3: VCS \u65e5\u5fd7\u8f93\u51fa",
              "\u5b9e\u9645\u8fd0\u884c\u5173\u952e\u65e5\u5fd7")
add_footer(slide, pg, TOTAL_PAGES)

add_code_block(slide, Inches(0.5), Inches(1.6), Inches(12.3), Inches(5.5),
               """[TB-VQ] DRIVER_OK detected, configuring VQ rings
[VQ] Queue 0 configured: desc=0x22c0000 avail=0x22c1000 used=0x22c1240 size=256
[VQ] Queue 1 configured: desc=0x1fb0000 avail=0x1fb1000 used=0x1fb1240 size=256
[TB-VQ] Virtqueue rings configured
[TB-VQ] RX queue notify (new buffers available)

[VQ-TX] desc[0]: addr=0x1954246 len=54 flags=0x0 next=1
[VQ-TX] Skipping 12 header bytes, reading 42 data bytes
[VQ-TX] Forwarded 42 bytes to ETH SHM (pkt #1)
[VQ-TX] Processed 1 descriptors, used_idx=1, total_tx=1
[TB-VQ] TX notify: processed 1 packets

[VQ-TX] desc[0]: addr=0x1954006 len=54 flags=0x0 next=1
[VQ-TX] Skipping 12 header bytes, reading 42 data bytes
[VQ-TX] Forwarded 42 bytes to ETH SHM (pkt #2)
[VQ-TX] Processed 1 descriptors, used_idx=2, total_tx=2
[TB-VQ] TX notify: processed 1 packets

[VQ-TX] desc[0]: addr=0x1954486 len=54 flags=0x0 next=1
[VQ-TX] Skipping 12 header bytes, reading 42 data bytes
[VQ-TX] Forwarded 42 bytes to ETH SHM (pkt #3)
[VQ-TX] Processed 1 descriptors, used_idx=3, total_tx=3
[TB-VQ] TX notify: processed 1 packets""")

# ======================================================================
# SECTION DIVIDER -- Part V: Phase 4
# ======================================================================
slide, pg = new_slide()
add_section_divider(slide, "V",
                    "\u7b2c\u4e94\u90e8\u5206: Phase 4 \u53cc\u5411\u7f51\u7edc\u901a\u4fe1",
                    "Phase 4: RX \u6ce8\u5165 + \u53cc VCS \u4e92\u8054 + Ping \u6d4b\u8bd5")
add_footer(slide, pg, TOTAL_PAGES)

# ======================================================================
# SLIDE 30 -- Phase 4: Architecture (Dual VCS Mode)
# ======================================================================
slide, pg = new_slide()
add_title_bar(slide, "Phase 4: \u53cc VCS \u4e92\u8054\u67b6\u6784",
              "\u4e24\u4e2a QEMU + \u4e24\u4e2a VCS \u901a\u8fc7 ETH SHM \u5b9e\u73b0\u53cc\u5411\u901a\u4fe1")
add_footer(slide, pg, TOTAL_PAGES)

tf = add_content_box(slide, Inches(0.5), Inches(1.6), Inches(5.5), Inches(5.5))
add_bullet(tf, "\u67b6\u6784\u8bbe\u8ba1", font_size=20, bold=True, color=DARK_BLUE)
add_bullet(tf, "Guest1 (10.0.0.2) + QEMU1 + VCS1 (Role A)", level=1, font_size=15)
add_bullet(tf, "Guest2 (10.0.0.1) + QEMU2 + VCS2 (Role B)", level=1, font_size=15)
add_bullet(tf, "ETH SHM \u5171\u4eab\u5185\u5b58\u4e92\u8054\u4e24\u4e2a VCS", level=1, font_size=15)
add_bullet(tf, "", font_size=8)
add_bullet(tf, "\u6570\u636e\u6d41", font_size=20, bold=True, color=DARK_BLUE)
add_bullet(tf, "TX: Guest \u2192 Virtqueue \u2192 VCS(TX) \u2192 ETH SHM", level=1, font_size=15)
add_bullet(tf, "RX: ETH SHM \u2192 VCS(RX) \u2192 Virtqueue \u2192 Guest", level=1, font_size=15)
add_bullet(tf, "", font_size=8)
add_bullet(tf, "\u4e24\u79cd\u6a21\u5f0f", font_size=20, bold=True, color=DARK_BLUE)
add_bullet(tf, "Mode 1 (TAP): \u5355 VCS + TAP bridge\uff08\u9700\u8981 sudo\uff09", level=1, font_size=15)
add_bullet(tf, "Mode 2 (DUAL): \u53cc VCS \u4e92\u8054\uff08\u65e0\u9700 sudo\uff09", level=1, font_size=15)

# Right side: architecture diagram
boxes_arch = [
    ("Guest1 (10.0.0.2)\nvirtio-net / eth0", Inches(7.0), Inches(1.6),
     RGBColor(0xE3, 0xF2, 0xFD)),
    ("QEMU1 + VCS1\n(Role A, MAC=01)", Inches(7.0), Inches(2.7), ACCENT_BLUE),
    ("ETH SHM\n(\u53cc\u5411\u73af\u5f62\u7f13\u51b2\u533a)", Inches(7.0), Inches(3.8),
     RGBColor(0xFC, 0xE4, 0xEC)),
    ("QEMU2 + VCS2\n(Role B, MAC=02)", Inches(7.0), Inches(4.9), ACCENT_BLUE),
    ("Guest2 (10.0.0.1)\nvirtio-net / eth0", Inches(7.0), Inches(6.0),
     RGBColor(0xE3, 0xF2, 0xFD)),
]
for txt, l, t, bg in boxes_arch:
    add_diagram_box(slide, l, t, Inches(5.5), Inches(0.85), txt, bg_color=bg,
                    font_size=13)

# Bidirectional arrows
for i in range(len(boxes_arch) - 1):
    mid_x = Inches(9.6)
    y_top = boxes_arch[i][2] + Inches(0.85)
    add_down_arrow(slide, mid_x, y_top, Inches(0.3), Inches(0.2))

# ======================================================================
# SLIDE 31 -- Phase 4: RX Injection Implementation
# ======================================================================
slide, pg = new_slide()
add_title_bar(slide, "Phase 4: RX \u6ce8\u5165\u5b9e\u73b0",
              "ETH SHM \u2192 VCS \u2192 Virtqueue \u2192 Guest")
add_footer(slide, pg, TOTAL_PAGES)

tf = add_content_box(slide, Inches(0.5), Inches(1.6), Inches(6), Inches(5.5))
add_bullet(tf, "RX \u6ce8\u5165\u6d41\u7a0b", font_size=20, bold=True, color=DARK_BLUE)
add_bullet(tf, "1. VCS \u4e3b\u5faa\u73af\u8f6e\u8be2 ETH SHM \u7684 RX \u961f\u5217", level=1, font_size=15)
add_bullet(tf, "2. \u6536\u5230\u5305\u540e DMA \u5199\u5165 Guest \u7684 RX Virtqueue \u7f13\u51b2\u533a", level=1, font_size=15)
add_bullet(tf, "3. \u6dfb\u52a0 virtio-net header (12B) \u524d\u7f00", level=1, font_size=15)
add_bullet(tf, "4. \u66f4\u65b0 used ring \u5e76\u89e6\u53d1 MSI/INTx \u4e2d\u65ad", level=1, font_size=15)
add_bullet(tf, "5. Guest \u5185\u6838\u6536\u5230\u4e2d\u65ad\u540e\u5904\u7406\u7f51\u7edc\u5305", level=1, font_size=15)
add_bullet(tf, "", font_size=8)
add_bullet(tf, "\u5173\u952e\u5b9e\u73b0\u7ec6\u8282", font_size=20, bold=True, color=DARK_BLUE)
add_bullet(tf, "virtqueue_rx_inject(): DMA \u5199 header + \u8d1f\u8f7d", level=1, font_size=15)
add_bullet(tf, "ISR \u5bc4\u5b58\u5668: BAR0+0x3000, \u8bfb\u6e05\u9664\u6a21\u5f0f", level=1, font_size=15)
add_bullet(tf, "INTx \u4e2d\u65ad: assert \u2192 Guest ISR \u8bfb \u2192 deassert", level=1, font_size=15)

add_code_block(slide, Inches(7.0), Inches(1.6), Inches(5.8), Inches(5.2),
"""// VCS \u4e3b\u5faa\u73af (tb_top.sv)
forever begin
  // 1. \u8f6e\u8be2 TLP (1ms \u8d85\u65f6, \u975e\u963b\u585e)
  ret = bridge_vcs_poll_tlp(...);
  if (ret == 0) begin
    // \u5904\u7406 TLP (MRd/MWr/CfgRd/CfgWr)
  end

  // 2. \u8f6e\u8be2 RX (\u4ece ETH SHM \u63a5\u6536)
  rx_len = eth_mac_rx_poll(rx_buf);
  if (rx_len > 0) begin
    virtqueue_rx_inject(rx_buf, rx_len);
    bridge_vcs_raise_msi(0);  // \u89e6\u53d1\u4e2d\u65ad
  end
end""")

# ======================================================================
# SLIDE 32 -- Phase 4: Problems & Fixes
# ======================================================================
slide, pg = new_slide()
add_title_bar(slide, "Phase 4: \u9047\u5230\u7684\u95ee\u9898\u4e0e\u89e3\u51b3")
add_footer(slide, pg, TOTAL_PAGES)

problems = [
    ["\u95ee\u9898", "\u539f\u56e0", "\u89e3\u51b3\u65b9\u6848"],
    ["TLP \u8f6e\u8be2\u963b\u585e\n\u5bfc\u81f4 RX \u65e0\u6cd5\u5de5\u4f5c",
     "sock_sync_recv \u963b\u585e\u7b49\u5f85\nVCS \u65e0\u6cd5\u8f6e\u8be2 ETH SHM",
     "\u65b0\u589e sock_sync_recv_timed()\n\u4f7f\u7528 poll() 1ms \u8d85\u65f6"],
    ["Guest2\u2192Guest1 Ping \u5931\u8d25\n\u5355\u5411\u901a\u4fe1\u6b63\u5e38",
     "Guest1 \u5b8c\u6210\u540e\u7acb\u5373\u5173\u673a\nVCS1 \u9000\u51fa\u65e0\u6cd5\u8f6c\u53d1",
     "\u589e\u52a0 30s \u5173\u673a\u524d\u7b49\u5f85\n\u786e\u4fdd\u53cc\u65b9\u6d4b\u8bd5\u5b8c\u6210"],
    ["ARP \u89e3\u6790\u5ef6\u8fdf\n\u5bfc\u81f4\u9996\u5305\u4e22\u5931",
     "\u53cc\u65b9\u540c\u65f6 Ping\n\u65e0 ARP \u7f13\u5b58",
     "\u589e\u52a0\u9884 Ping ARP \u63a2\u6d4b\u6b65\u9aa4\n\u786e\u4fdd MAC \u5730\u5740\u5df2\u7f13\u5b58"],
    ["VCS \u7f16\u8bd1\u5de5\u5177\u94fe\u95ee\u9898\n\u7cfb\u7edf GCC 4.8.5",
     "\u4e0d\u652f\u6301 C11 (stdatomic.h)\n\u94fe\u63a5\u5668\u4e0d\u652f\u6301 -no-pie",
     "\u4f7f\u7528 conda GCC 15.2.0\n-std=gnu11 + conda g++"],
]
add_table(slide, Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.5),
          len(problems), 3, problems,
          col_widths=[Inches(3.5), Inches(4.2), Inches(4.6)])

# ======================================================================
# SLIDE 33 -- Phase 4: Non-blocking Poll Fix
# ======================================================================
slide, pg = new_slide()
add_title_bar(slide, "Phase 4: \u6838\u5fc3\u4fee\u590d \u2014 \u975e\u963b\u585e TLP \u8f6e\u8be2",
              "\u89e3\u51b3 TX/RX \u8def\u5f84\u6b7b\u9501\u95ee\u9898")
add_footer(slide, pg, TOTAL_PAGES)

add_code_block(slide, Inches(0.5), Inches(1.5), Inches(6), Inches(2.5),
"""\u4fee\u590d\u524d (\u963b\u585e\u6a21\u5f0f):
int bridge_vcs_poll_tlp(...) {
    sync_msg_t msg;
    // \u963b\u585e\u7b49\u5f85 QEMU \u53d1\u9001 TLP
    int ret = sock_sync_recv(g_sock_fd, &msg);
    // VCS \u6c38\u8fdc\u963b\u585e\u5728\u8fd9\u91cc,
    // \u65e0\u6cd5\u8f6e\u8be2 ETH SHM \u63a5\u6536 RX \u5305!
}""")

add_code_block(slide, Inches(0.5), Inches(4.3), Inches(6), Inches(2.8),
"""\u4fee\u590d\u540e (\u975e\u963b\u585e\u6a21\u5f0f):
int bridge_vcs_poll_tlp(...) {
    sync_msg_t msg;
    // 1ms \u8d85\u65f6, \u975e\u963b\u585e
    int ret = sock_sync_recv_timed(
                  g_sock_fd, &msg, 1);
    if (ret == 1) return 1; // \u8d85\u65f6\u2192\u8f6e\u8be2RX
}""")

tf = add_content_box(slide, Inches(7.0), Inches(1.5), Inches(5.8), Inches(5.5))
add_bullet(tf, "\u95ee\u9898\u5206\u6790", font_size=20, bold=True, color=RED)
add_bullet(tf, "VCS \u4e3b\u5faa\u73af: poll_tlp \u2192 TLP\u5904\u7406 \u2192 RX\u8f6e\u8be2", level=1, font_size=15)
add_bullet(tf, "poll_tlp \u963b\u585e\u65f6, RX \u6c38\u8fdc\u4e0d\u4f1a\u88ab\u6267\u884c", level=1, font_size=15)
add_bullet(tf, "Guest \u53d1\u7684\u5305\u80fd\u51fa\u53bb(TX), \u4f46\u6536\u4e0d\u5230\u56de\u590d(RX)", level=1, font_size=15)
add_bullet(tf, "", font_size=8)
add_bullet(tf, "\u89e3\u51b3\u65b9\u6848", font_size=20, bold=True, color=GREEN)
add_bullet(tf, "\u65b0\u589e sock_sync_recv_timed() \u51fd\u6570", level=1, font_size=15)
add_bullet(tf, "\u4f7f\u7528 poll() \u7cfb\u7edf\u8c03\u7528, 1ms \u8d85\u65f6", level=1, font_size=15)
add_bullet(tf, "\u8d85\u65f6\u65f6\u8fd4\u56de 1, VCS \u7ee7\u7eed\u8f6e\u8be2 RX", level=1, font_size=15)
add_bullet(tf, "\u6709\u6570\u636e\u65f6\u6b63\u5e38\u5904\u7406 TLP", level=1, font_size=15)
add_bullet(tf, "", font_size=8)
add_bullet(tf, "\u6548\u679c", font_size=20, bold=True, color=GREEN)
add_bullet(tf, "TX \u548c RX \u8def\u5f84\u540c\u65f6\u5de5\u4f5c", level=1, font_size=15)
add_bullet(tf, "\u53cc\u5411 Ping \u5168\u90e8\u6210\u529f", level=1, font_size=15)

# ======================================================================
# SLIDE 34 -- Phase 4: Test Results
# ======================================================================
slide, pg = new_slide()
add_title_bar(slide, "Phase 4: \u6d4b\u8bd5\u7ed3\u679c",
              "\u53cc VCS \u4e92\u8054\u6a21\u5f0f \u2014 \u53cc\u5411 Ping \u6d4b\u8bd5")
add_footer(slide, pg, TOTAL_PAGES)

results = [
    ["\u6d4b\u8bd5\u9879", "\u65b9\u5411", "\u53d1\u9001", "\u63a5\u6536", "\u4e22\u5305\u7387", "\u72b6\u6001"],
    ["ARP \u63a2\u6d4b", "Guest1\u2192Guest2", "1", "1", "0%", "PASS"],
    ["ARP \u63a2\u6d4b", "Guest2\u2192Guest1", "1", "1", "0%", "PASS"],
    ["Ping \u6d4b\u8bd5", "Guest1\u2192Guest2", "3", "3", "0%", "PASS"],
    ["Ping \u6d4b\u8bd5", "Guest2\u2192Guest1", "3", "3", "0%", "PASS"],
]
add_table(slide, Inches(0.5), Inches(1.5), Inches(12.3), Inches(3.0),
          len(results), 6, results,
          col_widths=[Inches(2.5), Inches(2.8), Inches(1.2),
                      Inches(1.2), Inches(1.5), Inches(3.1)])

# Result summary
tf = add_content_box(slide, Inches(0.5), Inches(4.8), Inches(12.3), Inches(2.0))
p = add_bullet(tf, "Phase 4 \u7ed3\u679c: \u53cc\u5411 Ping ALL PASS (8/8 \u5305)",
               font_size=24, bold=True, color=GREEN)
add_bullet(tf, "", font_size=8)
add_bullet(tf, "\u6d4b\u8bd5\u914d\u7f6e", font_size=18, bold=True, color=DARK_BLUE)
add_bullet(tf, "Guest1: 10.0.0.2 (QEMU1 + VCS1/Role A, MAC=de:ad:be:ef:00:01)",
           level=1, font_size=14)
add_bullet(tf, "Guest2: 10.0.0.1 (QEMU2 + VCS2/Role B, MAC=de:ad:be:ef:00:02)",
           level=1, font_size=14)
add_bullet(tf, "ETH SHM: /cosim_eth0 (\u53cc\u5411\u73af\u5f62\u7f13\u51b2\u533a, a_to_b / b_to_a)",
           level=1, font_size=14)
add_bullet(tf, "\u4e2d\u65ad\u6a21\u5f0f: INTx (pci_set_irq assert/deassert)",
           level=1, font_size=14)

# ======================================================================
# SLIDE 35 -- Phase 4: VCS Log Output
# ======================================================================
slide, pg = new_slide()
add_title_bar(slide, "Phase 4: VCS \u65e5\u5fd7\u8f93\u51fa",
              "\u53cc\u65b9 TX/RX \u5305\u5904\u7406\u8bb0\u5f55")
add_footer(slide, pg, TOTAL_PAGES)

add_code_block(slide, Inches(0.3), Inches(1.5), Inches(6.3), Inches(5.5),
"""--- VCS1 (Role A) ---
[VQ-RX] Injected 98 bytes, rx_pkt #10
[VQ-RX] frame hex: de ad be ef 00 01
  de ad be ef 00 02 08 00 45 00 ...
  (ICMP echo reply from Guest2)
[EP-ISR] ISR read: value=0x00000001
[VQ-TX] desc[0]: addr=0x1954246 len=110
[VQ-TX] Skipping 12 header bytes
[VQ-TX] pkt_len=98 hex: de ad be ef 00 02
  de ad be ef 00 01 08 00 45 00 ...
  (ICMP echo reply to Guest2)
[VQ-TX] Forwarded 98 bytes to ETH SHM
  (pkt #10)""")

add_code_block(slide, Inches(6.8), Inches(1.5), Inches(6.3), Inches(5.5),
"""--- VCS2 (Role B) ---
[VQ-RX] Injected 98 bytes, rx_pkt #9
[VQ-RX] frame hex: de ad be ef 00 02
  de ad be ef 00 01 08 00 45 00 ...
  (ICMP echo request from Guest1)
[EP-ISR] ISR read: value=0x00000001
[VQ-TX] desc[0]: addr=0x1955206 len=110
[VQ-TX] Skipping 12 header bytes
[VQ-TX] pkt_len=98 hex: de ad be ef 00 01
  de ad be ef 00 02 08 00 45 00 ...
  (ICMP echo reply forwarded)
[VQ-TX] Forwarded 98 bytes to ETH SHM
  (pkt #10)""")

# ======================================================================
# SLIDE 36 -- Phase 5 Section Divider
# ======================================================================
slide, pg = new_slide()
add_section_divider(slide, 6, "Phase 5: TCP/iperf 吞吐测试",
                    "nc 多段 + iperf3 端到端吞吐量验证 — 全部 PASS")
add_footer(slide, pg, TOTAL_PAGES)

# ======================================================================
# SLIDE 37 -- Phase 5: Test Architecture
# ======================================================================
slide, pg = new_slide()
add_title_bar(slide, "Phase 5: 测试架构",
              "双 VCS 互联方案 — Server/Client 角色分配 + 静态 ARP")
add_footer(slide, pg, TOTAL_PAGES)

tf = add_content_box(slide, Inches(0.5), Inches(1.5), Inches(6), Inches(5.5))
add_bullet(tf, "测试目标", font_size=20, bold=True, color=DARK_BLUE)
add_bullet(tf, "验证 TCP 数据流端到端传输能力", level=1, font_size=16)
add_bullet(tf, "nc: 多段数据传输 (512B / 1KB / 2KB / 4KB)", level=1, font_size=16)
add_bullet(tf, "iperf3: 3s TCP 吞吐量基准测试", level=1, font_size=16)
add_bullet(tf, "", font_size=8)
add_bullet(tf, "双 VCS 互联方案", font_size=20, bold=True, color=DARK_BLUE)
add_bullet(tf, "两个 QEMU Guest 通过 ETH SHM 连接", level=1, font_size=16)
add_bullet(tf, "Server (Guest1): 10.0.0.1, MAC de:ad:be:ef:00:01",
           level=1, font_size=14, font_name=CODE_FONT)
add_bullet(tf, "Client (Guest2): 10.0.0.2, MAC de:ad:be:ef:00:02",
           level=1, font_size=14, font_name=CODE_FONT)
add_bullet(tf, "", font_size=8)
add_bullet(tf, "关键改进", font_size=20, bold=True, color=DARK_BLUE)
add_bullet(tf, "BQL 死锁修复: MSI 回调 + DMA 线程安全化", level=1, font_size=16,
           color=GREEN)
add_bullet(tf, "静态 ARP 表项: 防止测试间隙 ARP 缓存过期", level=1, font_size=16)

# Architecture diagram (right side)
tf2 = add_content_box(slide, Inches(6.8), Inches(1.5), Inches(6.2), Inches(5.5))
add_bullet(tf2, "数据流", font_size=20, bold=True, color=DARK_BLUE)
add_bullet(tf2, "", font_size=6)

add_code_block(slide, Inches(6.8), Inches(2.5), Inches(6.2), Inches(4.5),
"""Server (10.0.0.1)         Client (10.0.0.2)
  QEMU1 + VCS1(A)            QEMU2 + VCS2(B)
      |                           |
      |<-- ETH SHM (cosim_eth0) ->|
      |                           |
  [1] ping 连通测试           RTT ~7ms, 0% loss
      |                           |
  [2] nc -l -p 5000          nc -> 5000
      |  <-- 512B/1K/2K/4K --     |
      |       全部 PASS            |
  [3] iperf3 -s              iperf3 -c 10.0.0.1
      |  <-- 3s TCP 吞吐 ---      |
      |  9.74 Mbps sender         |
      |  8.98 Mbps receiver       |
      |  0 retransmissions        |

静态 ARP 映射:
  Server: 10.0.0.2 -> de:ad:be:ef:00:02
  Client: 10.0.0.1 -> de:ad:be:ef:00:01""")

# ======================================================================
# SLIDE 38 -- Phase 5: Test Results
# ======================================================================
slide, pg = new_slide()
add_title_bar(slide, "Phase 5: 测试结果",
              "所有测试全部 PASS — nc 多段 + iperf3 吞吐量验证")
add_footer(slide, pg, TOTAL_PAGES)

results = [
    ["测试项", "方向", "数据量", "结果", "耗时/吞吐", "状态"],
    ["Ping 连通", "双向", "ICMP", "0% 丢包", "RTT ~7ms", "PASS"],
    ["nc TCP 1KB", "Client→Server", "1024 bytes", "数据完整接收", "-", "PASS"],
    ["nc TCP 512B", "Client→Server", "512 bytes", "数据完整接收", "-", "PASS"],
    ["nc TCP 1024B", "Client→Server", "1024 bytes", "数据完整接收", "-", "PASS"],
    ["nc TCP 2048B", "Client→Server", "2048 bytes", "数据完整接收", "-", "PASS"],
    ["nc TCP 4096B", "Client→Server", "4096 bytes", "数据完整接收", "-", "PASS"],
    ["iperf3 3s", "Client→Server", "3s TCP", "9.74/8.98 Mbps", "0 retrans", "PASS"],
]
add_table(slide, Inches(0.5), Inches(1.5), Inches(12.3), Inches(4.0),
          len(results), 6, results,
          col_widths=[Inches(2.0), Inches(2.3), Inches(2.0),
                      Inches(2.5), Inches(1.5), Inches(2.0)])

# Result detail
tf = add_content_box(slide, Inches(0.5), Inches(5.8), Inches(6), Inches(1.5))
add_bullet(tf, "iperf3 吞吐量 (PASS)", font_size=20, bold=True, color=GREEN)
add_bullet(tf, "Sender: 9.74 Mbps", level=1, font_size=14, font_name=CODE_FONT)
add_bullet(tf, "Receiver: 8.98 Mbps", level=1, font_size=14, font_name=CODE_FONT)
add_bullet(tf, "Retransmissions: 0", level=1, font_size=14, font_name=CODE_FONT)

tf2 = add_content_box(slide, Inches(6.8), Inches(5.8), Inches(6.2), Inches(1.5))
add_bullet(tf2, "网络统计", font_size=20, bold=True, color=DARK_BLUE)
add_bullet(tf2, "Server: rx=2442 pkts, tx=366 pkts", level=1, font_size=14,
           font_name=CODE_FONT)
add_bullet(tf2, "Client: tx=2441 pkts, rx=366 pkts", level=1, font_size=14,
           font_name=CODE_FONT)
add_bullet(tf2, "0 errors, 0 dropped", level=1, font_size=14,
           font_name=CODE_FONT, color=GREEN)

# ======================================================================
# SLIDE 39 -- Phase 5: BQL 死锁修复 (关键技术突破)
# ======================================================================
slide, pg = new_slide()
add_title_bar(slide, "Phase 5: 关键技术突破 — BQL 死锁修复",
              "irq_poller 线程与主线程争锁问题的解决方案")
add_footer(slide, pg, TOTAL_PAGES)

# Problem description (left)
tf = add_content_box(slide, Inches(0.5), Inches(1.5), Inches(6), Inches(5.5))
add_bullet(tf, "问题: BQL 死锁", font_size=22, bold=True, color=RED)
add_bullet(tf, "", font_size=6)
add_bullet(tf, "irq_poller 线程在 MSI 回调中调用 bql_lock()",
           level=1, font_size=16)
add_bullet(tf, "主线程持有 BQL 等待 DMA 处理完成",
           level=1, font_size=16)
add_bullet(tf, "irq_poller 阻塞 → DMA 无法处理 → 死锁",
           level=1, font_size=16, color=RED)
add_bullet(tf, "", font_size=10)

add_code_block(slide, Inches(0.5), Inches(3.8), Inches(6), Inches(3.2),
"""死锁场景:
  主线程:    bql_lock() → 等待 DMA 完成
  irq_poller: MSI 回调 → bql_lock() → 阻塞!
              → DMA 回调无法执行
              → 主线程永远等待
              → 死锁!""")

# Solution (right)
tf2 = add_content_box(slide, Inches(6.8), Inches(1.5), Inches(6.2), Inches(2.0))
add_bullet(tf2, "解决方案", font_size=22, bold=True, color=GREEN)
add_bullet(tf2, "", font_size=6)
add_bullet(tf2, "MSI 回调: 改为入队 + QEMU bottom-half 延迟处理",
           level=1, font_size=16)
add_bullet(tf2, "irq_poller 永不阻塞 BQL", level=1, font_size=16, color=GREEN)
add_bullet(tf2, "DMA 回调: 改用 cpu_physical_memory_read/write",
           level=1, font_size=16)
add_bullet(tf2, "线程安全，无需 BQL", level=2, font_size=14, color=GREEN)

add_code_block(slide, Inches(6.8), Inches(3.8), Inches(6.2), Inches(3.2),
"""修复后:
  irq_poller: MSI → 入队 pending_msi
              → 不调用 bql_lock()
              → 立即返回

  主线程 BH: bql_lock() → 处理 pending_msi
             → MSI 注入 → bql_unlock()

  DMA 回调:  cpu_physical_memory_read/write()
             → 线程安全, 无需 BQL""")

# ======================================================================
# SLIDE 40 -- Phase 5: 网络统计与日志
# ======================================================================
slide, pg = new_slide()
add_title_bar(slide, "Phase 5: 网络统计与 VCS 日志",
              "Server/Client 包收发统计")
add_footer(slide, pg, TOTAL_PAGES)

# Network stats table
stats_data = [
    ["指标", "Server (10.0.0.1)", "Client (10.0.0.2)"],
    ["TX packets", "366", "2441"],
    ["RX packets", "2442", "366"],
    ["Errors", "0", "0"],
    ["Dropped", "0", "0"],
]
add_table(slide, Inches(0.5), Inches(1.5), Inches(8), Inches(2.8),
          len(stats_data), 3, stats_data,
          col_widths=[Inches(2.5), Inches(2.75), Inches(2.75)])

# VCS log excerpts
add_code_block(slide, Inches(0.3), Inches(4.8), Inches(6.3), Inches(2.5),
"""--- VCS1 (Role A / Server) ---
[VQ-TX] tx=366 packets forwarded to ETH SHM
[VQ-RX] rx=2442 packets injected from ETH SHM
  包含: ARP, ICMP, TCP (SYN/ACK/PSH/FIN)
  iperf3: 3s TCP 吞吐量测试数据流
[EP-ISR] MSI via bottom-half (无 BQL 争锁)
  0 errors, 0 dropped""")

add_code_block(slide, Inches(6.8), Inches(4.8), Inches(6.3), Inches(2.5),
"""--- VCS2 (Role B / Client) ---
[VQ-TX] tx=2441 packets forwarded to ETH SHM
[VQ-RX] rx=366 packets injected from ETH SHM
  包含: ARP, ICMP, TCP (SYN/ACK/PSH/FIN)
  iperf3: 9.74 Mbps sender, 8.98 Mbps receiver
[EP-ISR] MSI via bottom-half (无 BQL 争锁)
  0 errors, 0 dropped""")

# ======================================================================
# SLIDE 40 -- Key Technical Summary
# ======================================================================
slide, pg = new_slide()
add_title_bar(slide,
              "\u5173\u952e\u6280\u672f\u8981\u70b9\u603b\u7ed3")
add_footer(slide, pg, TOTAL_PAGES)

topics = [
    ("\u5b57\u8282\u5e8f\u5904\u7406 (LE MMIO)", [
        "PCIe MMIO \u901a\u8def\u4e3a Little-Endian",
        "QEMU mmio_ops \u8fd4\u56de host-endian",
        "VCS \u7aef\u6309 byte lane \u5bf9\u9f50",
        "MAC \u5730\u5740\u9700\u624b\u52a8 LE \u91cd\u6392",
    ], Inches(0.5), Inches(1.6)),
    ("Config Space \u5bf9\u9f50", [
        "CfgRd \u8fd4\u56de aligned dword",
        "byte_offset = reg & 0x3",
        "data >>= (byte_offset * 8)",
        "CfgWr \u540c\u7406\u9700\u8981 byte mask",
    ], Inches(6.8), Inches(1.6)),
    ("Virtqueue DMA \u5904\u7406", [
        "avail ring \u2192 desc chain \u2192 buffer DMA",
        "virtio-net header 12B \u5fc5\u987b\u5265\u79bb",
        "\u4efb\u610f\u957f\u5ea6 DMA: uint8_t* \u7f13\u51b2\u533a",
        "NBA \u65f6\u5e8f: notify \u9700\u7b49\u4e00\u4e2a clock",
    ], Inches(0.5), Inches(4.2)),
    ("\u975e\u963b\u585e\u8f6e\u8be2 + \u4e2d\u65ad\u914d\u7f6e", [
        "poll() 1ms \u8d85\u65f6\u907f\u514d TLP \u963b\u585e",
        "TX/RX \u53cc\u8def\u5f84\u4ea4\u66ff\u6267\u884c",
        "INTx: assert \u2192 ISR \u8bfb\u6e05 \u2192 deassert",
        "\u9759\u6001 ARP + TCP \u5355\u6bb5\u4f20\u8f93\u53ef\u9760",
    ], Inches(6.8), Inches(4.2)),
]

for title, items, left, top in topics:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(6), Inches(2.3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_GRAY
    shape.line.color.rgb = MEDIUM_BLUE
    shape.line.width = Pt(2)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.1)
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.font.name = CN_FONT

    for item in items:
        p2 = tf.add_paragraph()
        p2.text = f"  {item}"
        p2.font.size = Pt(14)
        p2.font.color.rgb = DARK_GRAY
        p2.font.name = CN_FONT
        p2.space_after = Pt(2)

# ======================================================================
# SLIDE 32 -- Summary & Outlook
# ======================================================================
slide, pg = new_slide()
add_title_bar(slide,
              "\u603b\u7ed3\u4e0e\u5c55\u671b")
add_footer(slide, pg, TOTAL_PAGES)

tf = add_content_box(slide, Inches(0.5), Inches(1.6), Inches(6), Inches(5))
add_bullet(tf,
           "\u5df2\u5b8c\u6210 (Phase 0 - 5)",
           font_size=20, bold=True, color=GREEN)
add_bullet(tf,
           "MMIO \u901a\u8def: MRd/MWr TLP \u5b8c\u6574\u8f6c\u53d1",
           level=1, font_size=16)
add_bullet(tf,
           "Config Space: CfgRd/CfgWr + BAR sizing",
           level=1, font_size=16)
add_bullet(tf,
           "Virtio \u5bc4\u5b58\u5668: common_cfg / notify / ISR / device_cfg",
           level=1, font_size=16)
add_bullet(tf,
           "\u9a71\u52a8\u7ed1\u5b9a: virtio_net probe \u6210\u529f, eth0 \u521b\u5efa",
           level=1, font_size=16)
add_bullet(tf,
           "DMA \u901a\u8def: \u53cc\u5411 DMA Read/Write \u9a8c\u8bc1",
           level=1, font_size=16)
add_bullet(tf,
           "Virtqueue TX/RX: ARP/ICMP/TCP \u5305\u8f6c\u53d1",
           level=1, font_size=16)
add_bullet(tf,
           "\u53cc VCS \u4e92\u8054: \u53cc\u5411 Ping 0% \u4e22\u5305",
           level=1, font_size=16)
add_bullet(tf,
           "TCP 吞吐测试: nc 多段 + iperf3 全部 PASS",
           level=1, font_size=16)
add_bullet(tf,
           "iperf3: 9.74 Mbps sender, 8.98 Mbps, 0 retrans",
           level=1, font_size=16)
add_bullet(tf,
           "BQL 死锁修复: MSI bottom-half + 线程安全 DMA",
           level=1, font_size=16)
add_bullet(tf, "", font_size=10)
add_bullet(tf,
           "关键技术突破",
           font_size=20, bold=True, color=GREEN)
add_bullet(tf,
           "BQL 死锁: irq_poller MSI 改为入队 + bottom-half",
           level=1, font_size=16)
add_bullet(tf,
           "DMA: cpu_physical_memory_read/write 线程安全",
           level=1, font_size=16)
add_bullet(tf,
           "全部 Phase 0-5 测试通过, 功能与性能均已验证",
           level=1, font_size=16)

phases = [
    ("Phase 0\nMMIO", GREEN, Inches(7.5)),
    ("Phase 1\nConfig", GREEN, Inches(7.5)),
    ("Phase 2\nVirtio", GREEN, Inches(7.5)),
    ("Phase 2b\nDriver", GREEN, Inches(7.5)),
    ("Phase 3\nVirtqueue", GREEN, Inches(7.5)),
    ("Phase 4\nPing", GREEN, Inches(7.5)),
    ("Phase 5\nTCP/iperf", GREEN, Inches(7.5)),
]
for i, (label, color, left) in enumerate(phases):
    y = Inches(1.6) + Inches(i * 0.75)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, y, Inches(5.3), Inches(0.58)
    )
    shape.fill.solid()
    if color == GREEN:
        shape.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)
    elif color == ORANGE:
        shape.fill.fore_color.rgb = RGBColor(0xFF, 0xF3, 0xE0)
    else:
        shape.fill.fore_color.rgb = RGBColor(0xE3, 0xF2, 0xFD)
    shape.line.color.rgb = color
    shape.line.width = Pt(2)

    tf2 = shape.text_frame
    p = tf2.paragraphs[0]
    if color == GREEN:
        status = "DONE"
    elif color == ORANGE:
        status = "PARTIAL"
    else:
        status = "PLANNED"
    p.text = f"{label}  --  {status}"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = color
    p.font.name = CN_FONT
    p.alignment = PP_ALIGN.CENTER


# ======================================================================
# Save
# ======================================================================
output_path = "/home/ubuntu/ryan/software/cosim_platform_report.pptx"
prs.save(output_path)
print(f"PPT saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
